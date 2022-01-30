# -*- coding: utf-8 -*-

"""
------------------------------
 MSDAC Systems - Participants
------------------------------

A basic participant organizer and generator for Seventh-day Adventist Church.

This module provides all functionality of organizing, customizing, exporting,
and automating the presentation layout of any church-related service.

Pre-requisites:
    - Microsoft(c) Office 2016 and above

Spacing Gap Formats:
    • Class           - 3 to 4 lines
    • Functions       - 2 to 3 lines
    • Block Comments  - 1 line
    • Inline Comments - varies

Disclaimer:
    This is an experimental program and we are aware of the slow performance
    due to Python's nature of interpreting code at runtime which affects speed.
    We are planning to switch to a better codebase later on though.


This program is part of MSDAC System's collection of softwares

Made with Qt, KV
(c) 2021-present Ken Verdadero, Reynald Ycong
"""


## Import Modules
import sys

try:
    import os, psutil, winreg, time, datetime, json, shutil, gc, re, bz2
    import Levenshtein as lev
    from PyQt5 import QtCore, QtGui, QtWidgets
    from PyQt5.QtCore import Qt
    from PyQt5.QtGui import QFont
    from pptx import Presentation
    from pptx.util import Inches, Cm, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from kenverdadero.KCore import KPath, KString
    from kenverdadero.KLogging import KLog
    from kenverdadero.KSoftware import KSoftware
    from kenverdadero.KCore.KCore import modHex, p, showLatency
except ImportError as e:
    print(f'Module not found: {e}')
    sys.exit()


        

class System(object):
    """
    System Class Handler
    
    Handles all system related functions for program to work including:
        - System variables
        - Verifying all required files and directory
        - Duplicate instances
        - Application exit event
    """
    def __init__(self):
        """
        All variables have the following format:
            self.TYPE_NAME_ATTRIBUTE_SUBATTR
            
            Example: self.SOFTWARE_HYMNALBROWSER_UI_SETTINGS
        """

        ## External File and Directories
        self.DIR_PARENT =           r'C:\ProgramData\MSDAC Systems'                                                         ## Parent Directory
        self.DIR_PROGRAM =          self.DIR_PARENT + r'\Participants'                                                      ## Program Directory
        self.DIR_LOG =              self.DIR_PROGRAM + r'\Logs'                                                             ## Log Directory
        self.FILE_DATA =            self.DIR_PROGRAM + r'\data.json'                                                        ## Both data and configuration are stored here
        self.FILE_PPT_EXPORTED =    self.DIR_PROGRAM + r'\exported.pptx'                                                    ## Path of the exported file

        ## Resources
        self.RES_LOGO =             './res/images/logo.png'                                                                 ## Resource SDA Logo Image
        self.RES_APP_ICON =         './res/images/app_icon.png'                                                             ## Resource Participants Main Icon
        self.RES_HEADERLOGO =       './res/images/header.png'                                                               ## Resource Main Header Image
        self.RES_DEFAULT_BG =       'res/images/defBG.png'                                                                  ## Resource Default Background Image
        self.RES_FONT_EXTS =        ['Cameliya.otf', 'HarrietTextBold.otf', 'HarrietTextBoldItalic.otf']                    ## Resource External Fonts for PowerPoint
        
        ## Properties
        self.PROCESS_NAME =         "participants.exe"                                                                      ## Program filename
        self.PROCESS =              psutil.Process(os.getpid())                                                             ## Get PID to detect multiple instances
        self.LOG_FILE_LIMIT =       10                                                                                      ## Maximum threshold for maintaining log files
        self.STARTUP_TIME =         0                                                                                       ## Set initial time for launching the program
        self.GLOBAL_STATE =         1                                                                                       ## 0 - Starting (unused), 1 - Ready, 2 - Reserved, 3 - Shutting Down
        self.EXT_MEMLIST =          "prt"                                                                                   ## Application Extension for Memberlist


    def verifyDirectories(self):
        """
        This method checks for directories and also generate new if the folders does not exist/
        This verifies from parent directory down to subfolders via loop using a dict of directories.

        ..  - Root
        ... - Etc
        <>  - File
        ->  - Directory

        Tree:
            .. MSDAC Systems (Parent/Root)
                -> Participants (Program)
                    -> Logs
                        <> someLogFile.log
                        <> ...
                    <> data.json

            -> POWERPNT.exe (MS Office)
        """
        DIRECTORIES = {
            self.DIR_PARENT: "Parent",
            self.DIR_PROGRAM: "Program",
            self.DIR_LOG: "Logs"
            }
        MISSING = 0

        for DIR, NAME in DIRECTORIES.items():
            if not KPath.exists(DIR, True):
                LOG.warn(f"{NAME} Directory \"{DIR}\" doesn't exist. Creating a new folder.")
                MISSING += 1
        

    def verifyRequisites(self):
        """
        Verifies all required programs to make the software run properly.
        This method prevents the program to run if a valid Microsoft Office
        PowerPoint is not found in both 32-bit and 64-bit installation folders
        """
        try:
            ## Uses Windows Registry to find the current path of MS Office
            self.PPT_EXEC = KPath.upFolder(winreg.EnumValue(winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                        r"SOFTWARE\Microsoft\Windows\CurrentVersion\App Paths\powerpnt.exe"), 0)[1])
            LOG.info(f'Using {"64" if "Program Files (x86)" not in self.PPT_EXEC else "32"}-Bit Version of Microsoft Office')
            LOG.info(f'Root Dir: {self.PPT_EXEC}')
        except FileNotFoundError:
            # No MS Office Installed/Detected
            LOG.crit("Microsoft Office is not installed or detected")
            MSG_BOX = QtWidgets.QMessageBox()
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Ok)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Critical)
            MSG_BOX.setText("Microsoft Office is not installed or detected.\nThis program uses Office PowerPoint to run properly.\n")
            MSG_BOX.setDetailedText("If you think this is a mistake, please contact the developers for further assistance.\n\nhttps://m.me/verdaderoken\nhttps://m.me/reynald.ycong")
            MSG_BOX.setWindowTitle(f"{SW.NAME} - Error")
            MSG_BOX.setWindowFlags(Qt.Drawer | Qt.WindowStaysOnTopHint)
            
            MSG_BOX.setStyleSheet(QSS.getStylesheet())
            MSG_BOX.exec_()
            LOG.sys("Program terminated due to an error: Cannot find PowerPoint directory.")
            sys.exit()
        return


    def checkInstances(self):
        """
        Asks the user if they want to run another instance
        """
        self.DUPLICATED = False
        DUPLICATES = list(filter(lambda x: x == SYS.PROCESS_NAME, [i.name() for i in psutil.process_iter()]))
        
        if len(DUPLICATES) > 2:
            MSG_BOX = QtWidgets.QMessageBox()
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Question)
            MSG_BOX.setText("The program is already running.\nDo you want to open another instance?")
            MSG_BOX.setWindowTitle("Duplicate Instance Detected")
            MSG_BOX.setWindowFlags(Qt.Drawer | Qt.WindowStaysOnTopHint)
            MSG_BOX.setStyleSheet(QSS.getStylesheet())
            MSG_BOX.setStyleSheet('min-width: 280px; min-height: 35px;')
            
            LOG.info("Duplicate Instance Detected")
            RET = MSG_BOX.exec_()
            if RET == QtWidgets.QMessageBox.Yes:
                return
            else:
                LOG.sys("Program terminated by user.")
                self.DUPLICATED = True
                sys.exit()
                
    
    def closeEvent(self, event):
        """
        Triggers this function when the UI is shut down
        Reports back to the logger.
        """
        
        if self.DUPLICATED:
            LOG.sys(f"Duplicate instance was shut down.")
        else:
            SYS.GLOBAL_STATE = 3
            LOG.sys(f"Shutting down")
            UIA.close()
            UIB.close()
            LOG.sys(f"Program terminated: Duration ({SW.runtime(2)} Seconds)")
            LOG.sys(f"END OF LOG - {datetime.datetime.now().strftime('%B %d, %Y - %I:%M:%S %p')}")




class Data(object):
    """
    Manages the pool database for Participants
    Uses JSON to parse the data.

    DCFG (or Database Configuration) is the dictionary of all recent role-name relations
    along with the configuration settings for users.
    """
    def __init__(self):
        self.check()


    def check(self):
        """
        Checks if the file is existent in directory
        """
        if not os.path.exists(SYS.FILE_DATA):
            LOG.warn("Data is missing. Generating new...")
            self.generateDefault()

        self.DATA = self.load()


    def load(self):
        """
        Retrieves the data from system's data file.
        """
        while True:
            with open(SYS.FILE_DATA, "r") as read:
                try: 
                    DATA = json.load(read)
                except json.decoder.JSONDecodeError:
                    LOG.crit('Failed to load statistic data. Regenerating default...')
                    self.generateDefault()
                else:
                    return DATA
        

    def dump(self, data=None, indent=4, sort_keys=True):
        """
        Saves the passed data to system's data file.
        Uses indention of 4 and sorted keys by default.
        Can be processed with other data if 2nd argument is specified.
        """
        if data is None: data = DCFG
        with open(SYS.FILE_DATA, "w") as write:
            json.dump(data, write, indent=indent, sort_keys=sort_keys)

    
    def generateDefault(self):
        """
        Generates default data
        """
        DATA = {
            "__DATECREATED__": time.time(),
            "__FILETYPE__": f"{SW.NAME} Data",
            "POOL": {
                "ROLES": ['Developed by'],
                "NAMES": ['MSDAC Systems']
            },
            "CONFIG": {
            }
        }
        self.dump(DATA)
        LOG.info(f"Default data was generated successfully. | Hash: {KString.toHashMD5(DATA)}")




class FileManager(object):
    """
    Manages all external files covered by the software
    This includes managing of recent files and temporary folder
    to maintain and prevent building up of unused data.
    """
    def __init__(self):
        self.deleteLogs()
        # self.checkExternalFonts()                                                                             ## Temporarily disabled


    def checkExternalFonts(self):
        """
        Auto installation of external custom fonts
        Needs administrator permissions.
        """
        FONTS = os.listdir("C:/Windows/Fonts")
        EXTS = ([f'res/fonts/{i}' for i in SYS.RES_FONT_EXTS])
        for f in EXTS:
            if f not in FONTS:
                shutil.copy(os.path.join(os.path.dirname(__file__), f), "C:/Windows/Fonts")


    def deleteLogs(self, deleteAll=False):
        """
        Orders the program to remove logs
        """
        self.deleteOldest(self.getLogFiles, SYS.LOG_FILE_LIMIT, deleteAll)


    def getLogFiles(self):
        """
        Returns list of log file paths
        """
        return [f'{SYS.DIR_LOG}\\{i}' for i in os.listdir(SYS.DIR_LOG) if i.endswith('.log')]


    def deleteOldest(self, fileList, threshold, deleteAll):
        """
        Delete older files when the folder items reached the maximum recent files allowed.
        """
        while len(fileList()) > (threshold if not deleteAll else 0):                                            ## Process code until the detected files are below threshold or when the switch is set to "Delete All" (0)
            try: os.remove(min(fileList(), key=os.path.getatime))                                               ## Eliminates the oldest accessed file
            except PermissionError as e: LOG.warn(e)                                                            ## Issue: There is no solution for this one yet. Administrator permissions could be used in future development
            except FileNotFoundError as e: pass                                                                 ## Ignore when file is not there anymore




class Stylesheet(object):
    """
    Handles all interface appearance for this program
    """
    def __init__(self):
        self.toggleMode()                                                                       ## Sets global palette for the application
        self.initStylesheet()                                                                   ## Sets the appearance of the UI's elements
    

    def setupFonts(self):
        """
        Initializes fonts to be used
        """
        FONT_BASE = QFont("Segoe UI", 9)
        FONT_BASE.setStyleStrategy(QFont.PreferAntialias)
        FONT_COLUMN = QFont('Segoe UI', 10)
        FONT_COLUMN.setStyleStrategy(QFont.PreferAntialias)
        FONT_COLUMN.setBold(True)
        APP.setFont(FONT_BASE)
        UIA.LBL_RL.setFont(FONT_COLUMN)
        UIA.LBL_NM.setFont(FONT_COLUMN)


    def initStylesheet(self):
        """
        Updates the stylesheet
        """
        SS = self.getStylesheet() 

        try:
            APP.setStyleSheet(SS)
            UIA.setStyleSheet(SS)
        except (NameError, AttributeError) as e:
            pass


    def QCl(self, c):
        """
        Returns QColor version of a hex value
        """
        if '#' in c: c = c[1:]
        return QtGui.QColor(int(c[0:2],16),int(c[2:4],16),int(c[4:6],16))
    

    def palette2Hex(self, color):
        """
        Returns HEX value of an RGB of a certain palette color
        """
        return self.RGBtoHEX(getattr(APP.palette(), color)().color().getRgb())
    
    
    def RGBtoHEX(self, rgb):
        rgb = list(rgb); del rgb[3]
        return '#%02x%02x%02x' % tuple(rgb)


    def toggleMode(self, mode=0):
        """
        Sets and updates all application color palette
        """

        if not mode:
            """
            For Light Mode
            """
            self.PRIMARY = '#004B74'
            self.SECONDARY = '#008A9A'
            self.TERTIARY = '#F7EBC5'
            self.GOLD = '#FFA92D'
            self.WARN = '#D25900'
            self.ERROR = '#9E1919'
            self.BORDER = '#C0C0C0'
            self.BORDER_HIGHLIGHT = '#C0C0C0'
            self.BTN_DISABLED = '#B6B6B6'
            self.TXT_INV = '#FFFFFF'
            self.TXT_DISABLED = '#999999'
            self.TXT_STATUSBAR = '#707070'
            self.STATUSBAR = '#CFCFCF'
            self.SCROLLBAR = '#A0A0A0'
            self.CARD = '#EEEEEE'
            self.CARDHOVER = '#EEEEEE'
            self.CTX_MENU = '#FFFFFF'
            PLT_LIGHT = QtGui.QPalette()
            PLT_LIGHT.setColor(QtGui.QPalette.Window, self.QCl('#FFFFFF'))
            PLT_LIGHT.setColor(QtGui.QPalette.WindowText, self.QCl('#202020'))
            PLT_LIGHT.setColor(QtGui.QPalette.Base, self.QCl('#DFDFDF'))
            PLT_LIGHT.setColor(QtGui.QPalette.AlternateBase, self.QCl('#2D2D2D'))
            PLT_LIGHT.setColor(QtGui.QPalette.ToolTipBase, self.QCl('#252525'))
            PLT_LIGHT.setColor(QtGui.QPalette.ToolTipText, self.QCl('#C5C5C5'))
            PLT_LIGHT.setColor(QtGui.QPalette.PlaceholderText, self.QCl('#999999'))
            PLT_LIGHT.setColor(QtGui.QPalette.HighlightedText, self.QCl('#EEEEEE'))
            PLT_LIGHT.setColor(QtGui.QPalette.Highlight, self.QCl(self.PRIMARY))
            PLT_LIGHT.setColor(QtGui.QPalette.Light, self.QCl('#D7D7D7'))
            PLT_LIGHT.setColor(QtGui.QPalette.Text, self.QCl('#202020'))
            PLT_LIGHT.setColor(QtGui.QPalette.Disabled, QtGui.QPalette.Text, self.QCl('#434343')) ## <- Unused
            PLT_LIGHT.setColor(QtGui.QPalette.Midlight, self.QCl('#888888'))
            PLT_LIGHT.setColor(QtGui.QPalette.Mid, self.QCl('#D2D2D2'))
            PLT_LIGHT.setColor(QtGui.QPalette.Dark, self.QCl('#555555'))
            PLT_LIGHT.setColor(QtGui.QPalette.Button, self.QCl('#CCCCCC'))
            PLT_LIGHT.setColor(QtGui.QPalette.Disabled, QtGui.QPalette.Button, self.QCl('#252525')) ## <- Unused
            PLT_LIGHT.setColor(QtGui.QPalette.ButtonText, self.QCl('#202020'))
            PLT_LIGHT.setColor(QtGui.QPalette.BrightText, self.QCl('#FFFFFF'))
            PLT_LIGHT.setColor(QtGui.QPalette.Link, self.QCl('#202020'))
            PLT_LIGHT.setColor(QtGui.QPalette.LinkVisited, self.QCl('#151515'))
            APP.setPalette(PLT_LIGHT)

        elif mode == 1:
            """
            For Dark Mode
            """
            self.PRIMARY = '#008A9A' ## #008A9A Original
            self.SECONDARY = '#004B74'
            self.TERTIARY = '#F7EBC5'
            self.GOLD = '#FFA92D'
            self.WARN = '#D25900'
            self.ERROR = '#9E1919'
            self.BORDER = '#303030'
            self.BORDER_HIGHLIGHT = '#505050'
            self.BTN_DISABLED = '#212121'
            self.TXT_INV = '#181818'
            self.TXT_DISABLED = '#434343'
            self.TXT_STATUSBAR = '#434343'
            self.STATUSBAR = '#2A2A2A'
            self.SCROLLBAR = '#3B3B3B'
            self.CARD = '#2A2A2A'
            self.CARDHOVER = '#323232'
            self.CTX_MENU = '#1D1D1D'
            PLT_DARK = QtGui.QPalette()
            PLT_DARK.setColor(QtGui.QPalette.Window, self.QCl('#202020'))
            PLT_DARK.setColor(QtGui.QPalette.WindowText, self.QCl('#D5D5D5'))
            PLT_DARK.setColor(QtGui.QPalette.Base, self.QCl('#191919'))
            PLT_DARK.setColor(QtGui.QPalette.AlternateBase, self.QCl('#2D2D2D'))
            PLT_DARK.setColor(QtGui.QPalette.ToolTipBase, self.QCl('#252525'))
            PLT_DARK.setColor(QtGui.QPalette.ToolTipText, self.QCl('#C5C5C5'))
            PLT_DARK.setColor(QtGui.QPalette.PlaceholderText, self.QCl('#999999'))
            PLT_DARK.setColor(QtGui.QPalette.HighlightedText, self.QCl('#191919'))
            PLT_DARK.setColor(QtGui.QPalette.Highlight, self.QCl(self.PRIMARY))
            PLT_DARK.setColor(QtGui.QPalette.Light, self.QCl('#898989'))
            PLT_DARK.setColor(QtGui.QPalette.Text, self.QCl('#EFEFEF'))
            PLT_DARK.setColor(QtGui.QPalette.Disabled, QtGui.QPalette.Text, self.QCl('#939393')) ## <- Unused
            PLT_DARK.setColor(QtGui.QPalette.Midlight, self.QCl('#888888'))
            PLT_DARK.setColor(QtGui.QPalette.Mid, self.QCl('#424242'))
            PLT_DARK.setColor(QtGui.QPalette.Dark, self.QCl('#555555'))
            PLT_DARK.setColor(QtGui.QPalette.Button, self.QCl('#353535'))
            PLT_DARK.setColor(QtGui.QPalette.Disabled, QtGui.QPalette.Button, self.QCl('#252525')) ## <- Unused
            PLT_DARK.setColor(QtGui.QPalette.ButtonText, self.QCl('#EFEFEF'))
            PLT_DARK.setColor(QtGui.QPalette.BrightText, self.QCl('#FFFFFF'))
            PLT_DARK.setColor(QtGui.QPalette.Link, self.QCl('#D0D0D0'))
            PLT_DARK.setColor(QtGui.QPalette.LinkVisited, self.QCl('#CECECE'))
            APP.setPalette(PLT_DARK)


    def getStylesheet(self, objectName=None):
        """
        Returns a string of stylesheet that will be used by QStyleSheet
        Values depends on what is the current theme.
        """
        # self.getThemes()
        RADIUS = "7px" ## Default: 9px
        RADIUS_SML = "4px" ## Default: 5px
        PADDING = "5px"

        if objectName is None:
            STYLESHEET = f"""
                QWidget#WIN_PARTICIPANTS {{
                    image: url('./res/images/bg.png');
                    image-position: bottom;
                }}



                /* Buttons */ 
                QPushButton {{
                    background-color: palette(button);
                    padding: {PADDING};
                    border-radius: {RADIUS};
                    border: 1px solid {modHex(self.palette2Hex('button'), 7)}
                }}
                QPushButton::pressed {{
                    background-color: palette(button);
                }}
                QPushButton::disabled {{
                    color: {self.TXT_DISABLED};
                    background-color: {self.BTN_DISABLED};
                }}
                QPushButton::hover {{
                    background-color: palette(light);
                }}
                /*
                QPushButton::focus {{
                    border: 1px solid {self.BORDER};
                }}
                */



                /* Export Powerpoint Button */ 
                QPushButton#BTN_POWERPOINT {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/ppt.png');
                }}

                QPushButton::hover#BTN_POWERPOINT {{
                    background-color: {self.PRIMARY};
                    image: url('./res/icons/ppt_hover.png');
                }}

                QPushButton::disabled#BTN_POWERPOINT {{
                    image: url('./res/icons/ppt_disabled.png');
                }}


                /* Export Plain Text Button */ 
                QPushButton#BTN_PLAINTEXT {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/plaintext.png');
                }}

                QPushButton::hover#BTN_PLAINTEXT {{
                    background-color: {self.PRIMARY};
                    image: url('./res/icons/plaintext_hover.png');
                }}

                QPushButton::disabled#BTN_PLAINTEXT {{
                    image: url('./res/icons/plaintext_disabled.png');
                }}



                /* Set Active Button  */ 
                QPushButton#BTN_ATVS {{
                    background-color: none;
                    border: none;
                    image: none;
                }}

                QPushButton::hover#BTN_ATVS {{
                    image: url('./res/icons/radio_translucent_hover.png');
                }}

                QPushButton::disabled#BTN_ATVS {{
                    image: url('./res/icons/radio_disabled.png');
                }}



                /* Set Active Button SELECTED */ 
                QPushButton#BTN_ATVS_SELECTED {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/radio_selected.png');
                }}

                QPushButton::hover#BTN_ATVS_SELECTED {{
                    image: url('./res/icons/radio_selected_hover.png');
                }}

                QPushButton::disabled#BTN_ATVS_SELECTED {{
                    image: url('./res/icons/select_radio_disabled.png');
                }}



                /* Edit Button */ 
                QPushButton#BTN_MEM_EDIT {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/edit.png');
                }}

                QPushButton::hover#BTN_MEM_EDIT {{
                    image: url('./res/icons/edit_hover.png');
                }}

                QPushButton::disabled#BTN_MEM_EDIT {{
                    image: url('./res/icons/edit_disabled.png');
                }}


                /* Insert/Add Button */ 
                QPushButton#BTN_INSS, QPushButton#BTN_MEM_ADD {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/add.png');
                }}

                QPushButton::hover#BTN_INSS, QPushButton::hover#BTN_MEM_ADD {{
                    image: url('./res/icons/add_hover.png');
                }}

                QPushButton::disabled#BTN_INSS, QPushButton::disabled#BTN_MEM_ADD  {{
                    image: url('./res/icons/add_disabled.png');
                }}



                /* Remove Button + Discard BG Img Button (Settings) + Remove */ 
                QPushButton#BTN_REMS, QPushButton#BTN_BG_DISCARD, QPushButton#BTN_MEM_REMOVE, QPushButton#BTN_GEN_REMOVE {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/xmark.png');
                }}

                QPushButton::hover#BTN_REMS, QPushButton::hover#BTN_BG_DISCARD, QPushButton::hover#BTN_MEM_REMOVE, QPushButton::hover#BTN_GEN_REMOVE {{
                    image: url('./res/icons/xmark_hover.png');
                }}

                QPushButton::disabled#BTN_REMS, QPushButton::disabled#BTN_BG_DISCARD, QPushButton::disabled#BTN_MEM_REMOVE, QPushButton::disabled#BTN_GEN_REMOVE {{
                    image: url('./res/icons/xmark_disabled.png');
                }}



                /* Import Export Button */
                QPushButton#BTN_MEM_IMPORT, QPushButton#BTN_GEN_IMPORT {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/import.png');
                }}

                QPushButton::hover#BTN_MEM_IMPORT, QPushButton::hover#BTN_GEN_IMPORT {{
                    image: url('./res/icons/import_hover.png');
                }}

                QPushButton::disabled#BTN_MEM_IMPORT, QPushButton::disabled#BTN_GEN_IMPORT {{
                    image: url('./res/icons/import_disabled.png');
                }}

                QPushButton#BTN_MEM_EXPORT, QPushButton#BTN_GEN_EXPORT {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/export.png');
                }}

                QPushButton::hover#BTN_MEM_EXPORT, QPushButton::hover#BTN_GEN_EXPORT {{
                    image: url('./res/icons/export_hover.png');
                }}

                QPushButton::disabled#BTN_MEM_EXPORT, QPushButton::disabled#BTN_GEN_EXPORT {{
                    image: url('./res/icons/export_disabled.png');
                }}



                /* Color Picker Button */
                QPushButton#BTN_FAC_COLORPICKER {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/fill.png');
                }}

                QPushButton::hover#BTN_FAC_COLORPICKER {{
                    image: url('./res/icons/fill_hover.png');
                }}

                QPushButton::disabled#BTN_FAC_COLORPICKER {{
                    image: url('./res/icons/fill_disabled.png');
                }}



                 /* Group Boxes */
                QGroupBox {{
                    border-radius: {RADIUS};
                    background-color: {self.CARD};
                    margin-top: 1.5em;
                    padding: 5px;
                    font-weight: bold;
                    font-size: 10pt;
                }}
                QGroupBox::hover {{
                    background-color: {self.CARDHOVER};
                }}
                QGroupBox::title {{
                    color: palette(text);
                    subcontrol-origin: margin;
                    left: 0px;
                    padding: 3px 5px 3px 5px;
                    border-radius: {RADIUS};
                }}
                QGroupBox::title::hover {{
                    border: 1px solid {self.PRIMARY};
                }}



                /* Remove Button LOCKED */ 
                QPushButton#BTN_REMS_LOCKED {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/locked.png');
                }}

                QPushButton::hover#BTN_REMS_LOCKED {{
                    image: url('./res/icons/unlock.png');
                }}

                QPushButton::disabled#BTN_REMS_LOCKED {{
                    image: url('./res/icons/locked_disabled.png');
                }}



                /* Save List Button */
                QPushButton#BTN_SAVELIST, QPushButton#BTN_MEM_SAVE {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/save.png');
                }}

                QPushButton::hover#BTN_SAVELIST, QPushButton::hover#BTN_MEM_SAVE {{
                    image: url('./res/icons/save_hover.png');
                }}

                QPushButton::disabled#BTN_SAVELIST, QPushButton::disabled#BTN_MEM_SAVE {{
                    image: url('./res/icons/save_disabled.png');
                }}



                /* Settings Button (Gear) */ 
                QPushButton#BTN_SETTINGS, QPushButton#BTN_GEN_MODIFY {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/settings.png');
                }}

                QPushButton::hover#BTN_SETTINGS, QPushButton::hover#BTN_GEN_MODIFY {{
                    image: url('./res/icons/settings_hover.png');
                }}


                /* Browse Button (Settings) */
                QPushButton#BTN_BG_BROWSE {{
                    background-color: none;
                    border: none;
                    image: url('./res/icons/folder.png');
                }}

                QPushButton::hover#BTN_BG_BROWSE {{
                    image: url('./res/icons/folder_hover.png');
                }}

                QPushButton::disabled#BTN_BG_BROWSE {{
                    image: url('./res/icons/folder_disabled.png');
                }}


                /* Dialog Boxes */
                QMessageBox {{
                    background-color: palette(window);
                }}
                


                /* Tooltip */
                QToolTip {{
                    color: palette(text);
                    background-color: palette(base);
                    border: none;
                }}



                /* Status Bar */
                QStatusBar#STATUSBAR {{
                    color: {self.TXT_STATUSBAR};
                    background-color: {self.STATUSBAR};
                }}



                /* Search Bars */
                QLineEdit, QComboBox{{
                    color: palette(text);
                    selection-color: {self.TXT_INV};
                    background-color: {self.CARD};
                    border: 1px solid {self.BORDER};
                    border-radius: {RADIUS};
                    padding: {PADDING};
                }}
                QLineEdit::focus#LNE_SEARCH {{
                    background-color: #AF{self.BORDER[1:]};
                }}
                QLineEdit::hover#LNE_SEARCH {{
                    border: 1px solid {self.BORDER_HIGHLIGHT};
                }}



                /* Combo Boxes */
                QComboBox {{
                    background-color: transparent;
                    border: 1px solid {self.BORDER};
                }}
                QComboBox::hover {{
                    background-color: transparent;
                    border: 1px solid {self.SECONDARY};
                }}

                QComboBox#CBX_RLS {{
                    margin-bottom: 1px;
                }}
                QComboBox#CBX_NMS {{
                    margin-bottom: 1px;
                }}



                /* Sliders */
                QSlider::handle {{
                    background-color: {self.PRIMARY};
                }}
                QSlider::handle::pressed {{
                    background-color: {self.SECONDARY};
                }}



                /* ScrollBars */
                QScrollBar:vertical {{
                    background-color: palette(base);
                    width: 16px;
                    margin: 0px;
                    border-radius: {RADIUS_SML};
                }}
                QScrollBar:horizontal {{
                    background-color: palette(base);
                    height: 15px;
                    margin: 0px;
                    border-radius: {RADIUS_SML};
                }}
                QScrollBar::handle:vertical {{
                    background-color: {self.SCROLLBAR}; min-height: 20px; margin: 3px; border-radius: {RADIUS_SML}; border: none;
                }}
                QScrollBar::handle:horizontal {{
                    background-color: {self.SCROLLBAR}; min-width: 20px; margin: 3px; border-radius: {RADIUS_SML}; border: none;
                }}
                QScrollBar::handle::hover {{
                    background-color: {modHex(self.SCROLLBAR, 20)}; min-width: 20px; margin: 3px; border-radius: {RADIUS_SML}; border: none;
                }}
                QScrollBar::handle::pressed {{
                    background-color: {self.PRIMARY}; min-width: 20px; margin: 3px; border-radius: {RADIUS_SML}; border: none;
                }}

                QScrollBar::add-line:vertical, QScrollBar::sub-line:vertical {{
                    border: none; background: none; height: 0px;
                }}
                QScrollBar::add-line:horizontal, QScrollBar::sub-line:horizontal {{
                    border: none; background: none; width: 0px;
                }}
                QScrollBar::up-arrow:vertical, QScrollBar::down-arrow:vertical, QScrollBar::add-page:vertical, QScrollBar::sub-page:vertical {{
                    border: none; background: none; color: none;
                }}
                QScrollBar::left-arrow:horizontal, QScrollBar::right-arrow:horizontal, QScrollBar::add-page:horizontal, QScrollBar::sub-page:horizontal {{
                    border: none; background: none; color: none;
                }}
                


                /* Settings Panel (List) */
                QAbstractItemView {{
                    color: palette(text);
                    outline: none;
                    background-color: palette(base);
                    border: none;
                    border-radius: {RADIUS};
                    selection-color: {self.TXT_INV};
                    selection-background-color: {self.SECONDARY};
                    padding: 3px;
                    min-height: 15px;
                }}
                QAbstractItemView::item {{
                    padding: 5px 2px 5px 4px;
                    margin: 2px 0px 2px 0px;
                    border-radius: {RADIUS};
                }}
                QAbstractItemView::item::selected {{
                    background-color: {self.PRIMARY};
                }}
                QAbstractItemView::item::hover {{
                    color: palette(text);
                    background-color: {self.CARD};
                }}
                QAbstractItemView::item::selected::hover {{
                    color: {self.TXT_INV};
                    background-color: qlineargradient(x1: 0, y1: 0, x2: 1, y2: 1,stop: 0 {modHex(self.PRIMARY, 50)}, stop: 1 {self.PRIMARY});
                }}

                /* Log Panel */
                QPlainTextEdit {{
                    color: palette(text);
                    background-color: palette(base);
                    border: 1px solid {self.BORDER};
                    border-radius: {RADIUS};
                }}
            


                /* Checkboxes (Toggle Switches) */
                QCheckBox {{
                    outline: none;
                }}
                QCheckBox::indicator {{
                    width: 23px;
                    height: 23px;
                }}
                QCheckBox::indicator::unchecked {{
                    image: url(./res/icons/off.png);
                }}
                QCheckBox::indicator::unchecked::hover {{
                    image: url(./res/icons/off_hover.png);
                }}
                QCheckBox::indicator::checked {{
                    image: url(./res/icons/on.png);
                }}
                QCheckBox::indicator::checked::hover {{
                    image: url(./res/icons/on_hover.png);
                }}
                QCheckBox::indicator::disabled {{
                    image: url(./res/icons/toggle_disabled.png);
                }}



                /* Spin

                /* Spin Box */
                QSpinBox {{
                    border: 1px solid {self.BORDER};
                    border-radius: {RADIUS};
                }}
                

                /* Add Queue Button */
                QPushButton#BTN_ADDQUEUE {{
                    background-color: transparent;
                    min-width: 17px;
                    width: 20px;
                    height: 20px;
                    image: url('./res/icons/add_.png');
                }}
                QPushButton::hover#BTN_ADDQUEUE {{
                    image: url('./res/icons/add_hover_.png');
                }}


                /* Queue Button */
                QPushButton#BTN_QUEUES {{
                    background-color: transparent;
                    min-width: 17px;
                    width: 20px;
                    height: 20px;
                    image: url('./res/icons/queue_.png');
                }}
                QPushButton::hover#BTN_QUEUES {{
                    image: url('./res/icons/queue_hover_.png');
                }}



                /* Special */

                QLabel#PIX_HEADER {{
                    color: {self.PRIMARY};
                }}
                QLabel::disabled, QLabel#LBL_BROWSERB {{
                    color: {self.TXT_DISABLED};
                }}
                
                QPushButton::enabled#BTN_LAUNCH, QPushButton::enabled#BTN_OK {{
                    color: {self.TXT_INV};
                    background-color: {self.PRIMARY};
                }}
                QPushButton#BTN_LAUNCH, QPushButton#BTN_OK  {{
                    min-width: 100px;
                }}
                QPushButton::hover#BTN_LAUNCH, QPushButton::hover#BTN_OK {{
                    color: {self.TXT_INV};
                    background-color: qlineargradient(x1: 0, y1: 0, x2: 1, y2: 1,stop: 0 {modHex(self.PRIMARY, 50)}, stop: 1 {self.PRIMARY});
                }}
                
                QPushButton::hover#BTN_RESET {{
                    color: palette(highlighted-text);
                    background-color: {self.ERROR};
                }}
                QLabel#LBL_RL, QLabel#LBL_NM {{
                    color: #005278;
                    padding: 3px;
                    border-radius: {RADIUS};
                }}
                """
            return STYLESHEET




class Package(object):
    """
    Represents a single package of configuration for a generated PowerPoint
    """
    def __init__(self):
        self.setupDefaultVariables()
        self.loadConfig()


    def restart(self):
        """
        Re-initiates all the process from the top
        """
        self.__init__()


    def setupDefaultVariables(self):
        """
        Initiates all default variables before loading the user-customized values
        Works as a placeholder incase something went wrong from user data (JSON)
        """
        self.DEF_DIR_EXPORT_RECENT = SW.DIR_CWD
        self.DIR_EXPORT_RECENT = SW.DIR_CWD
        self.DIR_EXPORT_MEMLIST = os.path.expanduser('~\Desktop')
        self.DIR_IMPORT_MEMLIST = os.path.expanduser('~\Desktop')
        self.DEF_IMG_BACKGROUND = "res/images/defBG.png"
        self.IMG_BACKGROUND = self.DEF_IMG_BACKGROUND
        self.FONT_TITLE = "Harriet Text Bold"
        self.FONT_SUBTITLE = "Cameliya"
        self.FONT_DATE = "Harriet Text Bold"
        self.FONT_CONTENT = "Harriet Text Bold"
        self.TXT_TITLE = "Sabbath Worship Participants"
        self.TXT_SUBTITLE = "Happy Sabbath!"
        self.RGB_TITLE = RGBColor(255, 255, 255)
        self.RGB_SUBTITLE = RGBColor(255, 169, 45)
        self.RGB_DATE = RGBColor(255, 255, 255)
        self.RGB_ROLES = RGBColor(221, 221, 221)
        self.RGB_NAMES = RGBColor(255, 255, 255)
        self.FSZ_TITLE = Pt(40)
        self.FSZ_SUBTITLE = Pt(30)
        self.FSZ_DATE = Pt(15)
        

    def loadConfig(self):
        """
        Retrieves all saved configurations from data.json
        Works by looping through all variables in this program and finding
        the working ones from absent ones.
        """
        LOG.info('Loading configuration')
        LOAD = [
            (self.IMG_BACKGROUND, "IMG_BACKGROUND"),
            (self.FONT_TITLE, "FONT_TITLE"),
            (self.FONT_SUBTITLE, "FONT_SUBTITLE"),
            (self.FONT_DATE, "FONT_DATE"),
            (self.FONT_CONTENT, "FONT_CONTENT"), 
            (self.TXT_TITLE, "TXT_TITLE"),
            (self.TXT_SUBTITLE, "TXT_SUBTITLE"),
            (self.DIR_EXPORT_RECENT, "DIR_EXPORT_RECENT"),
            (self.DIR_IMPORT_MEMLIST, "DIR_IMPORT_MEMLIST"),
            (self.DIR_EXPORT_MEMLIST, "DIR_EXPORT_MEMLIST")
        ]
        for obj in LOAD:
            try:
                for var in dict(self.__dict__):
                    if getattr(self, var) == obj[0]:
                        setattr(self, obj[1], DCFG["CONFIG"][obj[1]])
                        break
            except Exception as e:
                pass

        LOG.info('Successfully loaded configuration.')




class Core(object):
    """
    Contains methods that are used for better user experience
    """
    def __init__(self):
        self.IGNORED_CASE = ['of', 'to', 'and', 'by', 'for']


    def centerWindow(self, ui):
        """
        Relocates the specific UI from argument to center of the screen
        """
        FRM_GEOMETRY = ui.frameGeometry()
        SCREEN = QtWidgets.QApplication.desktop().screenNumber(QtWidgets.QApplication.desktop().cursor().pos())
        FRM_GEOMETRY.moveCenter(QtWidgets.QApplication.desktop().screenGeometry(SCREEN).center())
        ui.move(FRM_GEOMETRY.topLeft())


    def centerInsideWindow(self, WGTA, WGTB):
        """
        Relocates the Widget A centered with relation to Widget B
        """
        WINDOW = (WGTB.frameGeometry().width()-WGTA.frameGeometry().width(),
                WGTB.frameGeometry().height()-WGTA.frameGeometry().height())
        WGTA.move(WGTB.pos().x()+int(WINDOW[0]/2), WGTB.pos().y()+int(WINDOW[1]/2.2))


    def splitContents(self, roles:str, names:str):
        """
        Splits the content to distinguish Sabbath school from Divine Service.
        Has fallback for when the service-based categorizing fails
        """
        for i, cmb in enumerate(FLD.CBX_RLS):
            ## Regular Splitter
            if cmb.currentText().upper() in ["CLOSING PRAYER", "CLOSINGPRAYER", "CLOSING"]:
                SS_A = roles.split('\n', i+1); DS_A = SS_A[-1].split('\n'); SS_A.pop(-1)
                SS_B = names.split('\n', i+1); DS_B = SS_B[-1].split('\n'); SS_B.pop(-1)
                return (SS_A, SS_B, DS_A, DS_B)
                
        ## Basic Splitter (Fallback)
        SS_A = roles.split('\n', int(len(roles.split('\n'))/2)); DS_A = SS_A[-1].split('\n', int(len(SS_A[-1].split('\n'))/2)); SS_A.pop(-1)
        SS_B = names.split('\n', int(len(names.split('\n'))/2)); DS_B = SS_B[-1].split('\n', int(len(SS_B[-1].split('\n'))/2)); SS_B.pop(-1)
        return (SS_A, SS_B, DS_A, DS_B)


    def adjustRoleFormat(self, role:str):
        """
        Return the role string to a presentable format
        """
        return ' '.join(map(str, [f"{w[:1].upper()}{w[1:].lower()}" if w.lower() not in
                self.IGNORED_CASE else w.lower() for w in role.split()]))




class Fields(object):
    """
    Handles all field-related events:
        - Insert/Duplicate field(s)
        - Remove field(s)
        - Lock/Unlock fields
        - Button refreshers
        - Auto-add items to combo box after editing
        - Field data access
    
    Field - A group of widgets linked to one participant.
            It contains role and name combo boxes and the operator buttons (active, insert/duplicate, remove/lock)
    """
    def __init__(self):
        """
        Constructor
        """
        self.LOADING = True
        self.CBX_RLS, self.CBX_NMS, self.BTN_REMS, self.BTN_INSS, self.BTN_ATVS = [], [], [], [], []
        self.FIELDS = 0
        self.FIELDS_MAX = 20
        self.PREV_ACTIVE = None
        self.PREV_CBX = (0,0,0,0)
        self.MAX_VISIBLE_ITEMS = 13
        self.SEPARATOR = "—————————"

        ## Tooltip Values
        self.TTIP_BTN_ATVS = "Set as active"
        self.TTIP_BTN_ATVS_SELECTED = "Unset from being active"
        self.TTIP_BTN_INSS = "Add new field below or right click to duplicate"
        self.TTIP_BTN_REMS = "Remove the field or right click to lock this field"
        self.TTIP_BTN_REMS_LOCKED = "Right click to unlock"


    def setup(self):
        """
        Initiates the first run of field generation for setup.
        Usually triggered after setting up the UI
        """
        LENGTH = len(DCFG['POOL']['ROLES'])
        if not LENGTH:                                                                                  ## Eliminates no-field issue when there's no data from pool
            PDB.generateDefault()
            LENGTH = 1
        self.addFields(LENGTH)
        self.fillupItems()
        self.refreshItems()
        self.refreshStates()
        self.LOADING = False


    def refreshStates(self):
        """
        Refreshes the states of every field.

        Updates the connections of every button.
        This method helps the class handler to reconnect all connections
        of every present button after recent changes
        """
        ## Reconnect Signals
        for i in range(len(self.BTN_REMS)):                                                             ## Loops through every single button object based on BTN_REMS or BTN_INSS
            try:
                self.CBX_RLS[i].lineEdit().editingFinished.disconnect()
                self.CBX_NMS[i].lineEdit().editingFinished.disconnect()
                self.CBX_RLS[i].lineEdit().textChanged.disconnect()
                self.CBX_NMS[i].lineEdit().textChanged.disconnect()
                self.BTN_ATVS[i].clicked.disconnect()
            except Exception:                                                                           ## Throws exception when the buttons aren't connected yet which is not valid in the disconnect method
                pass
            finally:                                                                                    ## Always connect the Add and Remove button to its main method
                self.CBX_RLS[i].lineEdit().editingFinished.connect(lambda: self.refreshItems())
                self.CBX_NMS[i].lineEdit().editingFinished.connect(lambda: self.refreshItems())
                self.CBX_RLS[i].lineEdit().textChanged.connect(lambda: self.recordCbx('RLS'))
                self.CBX_NMS[i].lineEdit().textChanged.connect(lambda: self.recordCbx('NMS'))
                self.BTN_ATVS[i].clicked.connect(lambda: self.setActiveField())
                self.BTN_INSS[i].mouseReleaseEvent = lambda event: self.mouseReleased('INSS', event)
                self.BTN_REMS[i].mouseReleaseEvent = lambda event: self.mouseReleased('REMS', event)
        
        ## Prevent Overflow
        for btn in self.BTN_INSS:
            if self.FIELDS >= self.FIELDS_MAX:
                btn.setEnabled(False)                                                                   ## Sets the button to whether Enabled or Disabled depending on how many fields are active
                btn.setToolTip("")
            else:
                btn.setEnabled(True)
                btn.setToolTip(self.TTIP_BTN_INSS)

        ## Prevent Depletion
        try:
            self.BTN_REMS[0].setEnabled(False if self.FIELDS <= 1 else True)                            ## Disables the last field's remove button
        except IndexError:
            pass
        
        UIA.updateWindowTitle()


    def refreshItems(self):
        """
        Overrides the LineEdit's editingFinished event
        Adds the edited text to an item to every combo boxes if eligible
        """
        m = time.time()
        PRCD = (self.CBX_RLS, RLS), (self.CBX_NMS, NMS)                                 ## Procedure Variable (Role Objects, Role List) & (Name Objects, Name List)
        s = 0                                                                           ## Step
        
        try:
            HCBX = self.PREV_CBX                                                        ## Used for holding values for previous recorded CBX
            CMB_RECENT = PRCD[HCBX[0]][0][HCBX[1]]
            CONVERTED = (CORE.adjustRoleFormat(CMB_RECENT.currentText())                ## Converts the string to name-case
                if not HCBX[0] else MEM.adjustNameFormat(CMB_RECENT.currentText()))
            CMB_RECENT.setCurrentText(CONVERTED)
        except IndexError:
            pass

        while True:
            ITEMS = [c.currentText().strip() for c in PRCD[s][0]]                       ## Retrieve all current items displayed

            for i, c in enumerate(PRCD[s][0]):                                          ## Clear and fill up all CBX objects with items
                c.clear()
                c.addItems(ITEMS)
                c.setCurrentIndex(i)
                ## FIX THIS (Blank spaces)
                if c.currentText() == '':                                               ## Prevent including the blank item in combo box by resetting the index to none (-1)
                    c.setCurrentIndex(-1)
                    c.removeItem(i)
            
            MERGE = [item for item in PRCD[s][1] if item not in ITEMS]                  ## Merge with excess names from pool
            if len(MERGE):
                for i, c in enumerate(PRCD[s][0]):
                    # c.insertSeparator(c.count()-1)  ## <-- Not working for some reason
                    c.addItem(self.SEPARATOR)
                    c.addItems(MERGE)


            # ## Scan for identicals
            # PARTIAL = ITEMS + MERGE
            # IDENTICAL = [i for i in range(len(PRCD[s][0])) if PRCD[s][0][i].itemText(i) == HCBX[3]]
            # IDENTICAL_TEXT = list(set([PRCD[s][0][i].itemText(i) for i in range(len(PRCD[s][0])) if PRCD[s][0][i].itemText(i) == HCBX[3]]))


            # if HCBX[3] != 0:
            #     p(f"HOLD 3: {HCBX[3]}")
            #     p(f"IDENTICAL: {IDENTICAL[1:]}")
            #     p(f"IDENTICAL_TEXT: {IDENTICAL_TEXT[1:]}")
            #     duplicated = False                                                          ## Indicator; 
            #     for i, c in enumerate(PRCD[s][0]):                                          ## Clear and fill up all CBX objects with items
            #         c.clear()
            #         placed = False                                                          ## Indicator; when set to False, the duplicated text is not yet added.

            #         for j, item in enumerate(PARTIAL):                                      ## Add Items except the duplicated ones (allowing only 1 entry for duplicates)
            #             if item in IDENTICAL_TEXT:
            #                 if placed: continue
            #                 else: placed = True
            #             # if item == '': continue
            #             # p(f"Adding {item} {j} {placed}")
            #             if j == len(ITEMS): c.addItem(self.SEPARATOR)
            #             c.addItem(item)

            #         # p(f'Current Index: {i} {c.currentText()}')
            #         p([c.itemText(x) for x in range(c.count())])
            #         if i in IDENTICAL[1:]:
            #             p(f'Rearranging duplicate index: {i} {c.currentText()}')
            #             c.setCurrentIndex(IDENTICAL[0])
            #             duplicated = True
            #         else:
            #             # p(f'Rearranging normal index: {i} {c.currentText()}')
            #             play = i if not duplicated else i-(len(IDENTICAL)-1)
            #             # p(play, duplicated)
            #             c.setCurrentIndex(play)
            #         # for l in range(self.FIELDS):
            #         #     c.removeItem(0)

            ## NOTE: 
            ##      should remove blank spaces in items for every empty field
            ##

            if s > 0: break
            else: s += 1


    def fillupItems(self):
        """
        Syncs the entries for every field.

        This method also removes blank string retrieved from a pool.
        """
        for i,(r,n) in enumerate(zip(self.CBX_RLS, self.CBX_NMS)):
            r.setCurrentIndex(i)                                                            ## Sets the index to its exact position
            n.setCurrentIndex(i)

            if r.currentText() == '': r.setCurrentIndex(-1)                                 ## Rebind to none
            if n.currentText() == '': n.setCurrentIndex(-1)


    def redirectFieldInsertion(self, duplicate=False):
        """
        Redirects function to UIA's Field Insertion function.
        Triggers from the Add button via PyQt signal 
        """
        m = time.time()
        for i, btn in enumerate(self.BTN_INSS):
            if not btn.hasFocus(): continue
            self.insertField(i)
            self.refreshItems()
            if duplicate:
                self.CBX_RLS[i+1].setCurrentIndex(self.CBX_RLS[i].findText(self.CBX_RLS[i].currentText()))
                self.CBX_NMS[i+1].setCurrentIndex(self.CBX_NMS[i].findText(self.CBX_NMS[i].currentText()))
            break
        showLatency(m)


    def setActiveField(self):
        """
        Sets an active field by exporting its current values into a text file 
        to be read by a Text Source from OBS Studio
        """
        for i, btn in enumerate(self.BTN_ATVS):
            if not btn.hasFocus(): continue
            if btn.objectName() == 'BTN_ATVS_SELECTED':                                             ## Unset the field from being active
                EXP.fromActiveField(i, True)
                btn.setObjectName('BTN_ATVS')
            else:                                                                                   ## Before setting the triggered button to active, determine the previous
                if self.PREV_ACTIVE is not None:                                                    ## index to unset from being active and return into a normal state.
                    POINTER = self.PREV_ACTIVE
                    try:
                        self.BTN_ATVS[POINTER]
                    except IndexError:
                        POINTER = len(self.BTN_ATVS)-1
                    finally:
                        self.BTN_ATVS[POINTER].setObjectName('BTN_ATVS')
                        self.BTN_ATVS[POINTER].setToolTip(self.TTIP_BTN_ATVS)
                        self.BTN_ATVS[POINTER].setStyleSheet(QSS.getStylesheet())
                
                EXP.fromActiveField(i)                                                              ## Set the focused button to be active and export the file
                btn.setObjectName('BTN_ATVS_SELECTED')
                btn.setToolTip(self.TTIP_BTN_ATVS_SELECTED)
                self.PREV_ACTIVE = i

            btn.setStyleSheet(QSS.getStylesheet())
            break


    def insertField(self, pos:int=-1):
        """
        Inserts one new set of widgets with items available from both pools (roles, names).
        Creates Role, Name, Add button, Remove button to be placed on a vertical box layout.

        QVBoxLayout was used because of its insert method which you can use for inserting new
        widget to a specific index position.

        There are some adjustments for `pos` argument because of the nature of how Qt handles
        where it puts the next item in layout.

        `CBX_RLS` (Combo Box Roles) - contains all objects for roles
        `CBX_NMS` (Combo Box Names) - contains all objects for names
        `BTN_REMS` (Remove Push Button(s)) - contains all objects for remove/clear button
        `BTN_INSS` (Insert Push Button(s)) - contains all objects for insert/plus/add button
        `pos` - a positional marker used to mark where the next item should be placed

        Code block below will execute the following in order:
        1. Determining the right positional marker to prepare for insertion
        2. Inserting a placeholder (NoneType) in the newly placed index
        3. Compensates `pos` for zero-based indexing
        4. Starts instantiating all four (4) widgets to make a single-line field
        5. Combo boxes are filled with data from the pool for both Roles and Names
        6. Inserting the item in a QVBoxLayout
        7. Refresh states

        This method is not a loop but it can be used multiple times using a different method
        `addFields` which accepts numbers of how many times you want to generate a field.
        """
        # p([str(i)[-6:-1] for i in self.BTN_REMS]) ## Addresses (For debugging)

        ## Index-related adjustments
        pos +=1                                                                                 ## Increase position number by 1 to prevent placing the object in the first index
        if pos == len(self.CBX_RLS): pos += 1                                                   ## When last button is clicked, increment pos by 1 to ensure that the placeholder will be at the very last
        elif pos == 0: pos -= 1                                                                 ## When the pos is at the start, decrease the pos to compensate with insert function
            
        ## Appends a placeholder for an object to be created
        self.CBX_RLS.insert(pos, None)
        self.CBX_NMS.insert(pos, None)
        self.BTN_REMS.insert(pos, None)
        self.BTN_ATVS.insert(pos, None)
        self.BTN_INSS.insert(pos, None)

        if pos == len(self.CBX_RLS): pos -= 1                                                   ## When the pos is at the last index, decrease by 1 to compensate with zero-indexing
        FDS = len(self.CBX_RLS)-1                                                               ## Retrieve number of fields present in the layout
        self.FIELDS = FDS+1

        ## Fills the placeholder with objects (for Role, Name, Add, and Clear/Remove button) 
        self.CBX_RLS[pos] = QtWidgets.QComboBox(UIA.WGT_CENTRAL); self.CBX_RLS[pos].setObjectName(f"CBX_RLS")
        self.CBX_RLS[pos].view().window().setWindowFlags(Qt.Popup | Qt.FramelessWindowHint | Qt.NoDropShadowWindowHint)
        self.CBX_RLS[pos].view().window().setAttribute(Qt.WA_TranslucentBackground)
        self.CBX_RLS[pos].setMinimumWidth(140); self.CBX_RLS[pos].setMaximumWidth(200)
        self.CBX_RLS[pos].setEditable(True)
        self.CBX_RLS[pos].addItems(RLS)
        self.CBX_RLS[pos].setCurrentText('')
        self.CBX_RLS[pos].setMaxVisibleItems(self.MAX_VISIBLE_ITEMS)
        self.CBX_RLS[pos].wheelEvent = lambda e: self.ignoreWheel(e)

        self.CBX_NMS[pos] = QtWidgets.QComboBox(UIA.WGT_CENTRAL); self.CBX_NMS[pos].setObjectName(f"CBX_NMS")
        self.CBX_NMS[pos].view().window().setWindowFlags(Qt.Popup | Qt.FramelessWindowHint | Qt.NoDropShadowWindowHint)
        self.CBX_NMS[pos].view().window().setAttribute(Qt.WA_TranslucentBackground)
        self.CBX_NMS[pos].setMinimumWidth(140); self.CBX_NMS[pos].setMaximumWidth(200)
        self.CBX_NMS[pos].setEditable(True)
        self.CBX_NMS[pos].addItems(NMS)
        self.CBX_NMS[pos].setCurrentText('')
        self.CBX_NMS[pos].setMaxVisibleItems(self.MAX_VISIBLE_ITEMS)
        self.CBX_NMS[pos].wheelEvent = lambda e: self.ignoreWheel(e)

        self.BTN_ATVS[pos] = QtWidgets.QPushButton(UIA.WGT_CENTRAL); self.BTN_ATVS[pos].setObjectName(f"BTN_ATVS")
        self.BTN_ATVS[pos].setMaximumWidth(25);        self.BTN_ATVS[pos].setMinimumHeight(30)
        self.BTN_ATVS[pos].setToolTip(self.TTIP_BTN_ATVS)
        self.BTN_ATVS[pos].setFocusPolicy(Qt.ClickFocus)

        self.BTN_INSS[pos] = QtWidgets.QPushButton(UIA.WGT_CENTRAL); self.BTN_INSS[pos].setObjectName(f"BTN_INSS")
        self.BTN_INSS[pos].setMaximumWidth(26); self.BTN_INSS[pos].setMinimumHeight(30)
        self.BTN_INSS[pos].setToolTip(self.TTIP_BTN_INSS)
        self.BTN_INSS[pos].setFocusPolicy(Qt.ClickFocus)

        self.BTN_REMS[pos] = QtWidgets.QPushButton(UIA.WGT_CENTRAL); self.BTN_REMS[pos].setObjectName(f"BTN_REMS")
        self.BTN_REMS[pos].setMaximumWidth(25); self.BTN_REMS[pos].setMinimumHeight(30)
        self.BTN_REMS[pos].setToolTip(self.TTIP_BTN_REMS)
        self.BTN_REMS[pos].setFocusPolicy(Qt.ClickFocus)

        ## Finally adds those widgets into vertical layouts
        UIA.LYT_ROLES.insertWidget(pos, self.CBX_RLS[pos])
        UIA.LYT_NAMES.insertWidget(pos, self.CBX_NMS[pos])
        UIA.LYT_ACTIV.insertWidget(pos, self.BTN_ATVS[pos])
        UIA.LYT_CLEAR.insertWidget(pos, self.BTN_REMS[pos])
        UIA.LYT_INSRT.insertWidget(pos, self.BTN_INSS[pos])

        self.refreshStates()
    

    def addFields(self, fields:int=2):
        """
        Allows to add more field(s) using 2nd argument and redirects to `insertField` method
        """
        for i in range(fields): self.insertField(i)


    def removeField(self):
        """
        Removes a specific field by unlinking the widget from
        the layout and deleting its widgets from memory.
        """
        for i, btn in enumerate(self.BTN_REMS):
            if not btn.hasFocus(): continue
            if btn.objectName() != "BTN_REMS": return

            UIA.LYT_ROLES.removeWidget(self.CBX_RLS[i])
            UIA.LYT_NAMES.removeWidget(self.CBX_NMS[i])
            UIA.LYT_INSRT.removeWidget(self.BTN_ATVS[i])
            UIA.LYT_INSRT.removeWidget(self.BTN_INSS[i])
            UIA.LYT_CLEAR.removeWidget(self.BTN_REMS[i])

            del self.CBX_RLS[i]
            del self.CBX_NMS[i]
            del self.BTN_REMS[i]
            del self.BTN_INSS[i]
            del self.BTN_ATVS[i]
            self.FIELDS -= 1
            self.refreshStates()
            gc.collect()

            ## Reset window to shortest possible to remove spaces left by the field.
            UIA.resize(UIA.size().width(), 0)   
            break
    

    def lockUnlockField(self):
        """
        Handles lock and unlock mechanism of the field selected
        """
        for i, btn in enumerate(self.BTN_REMS):
            if not btn.hasFocus(): continue
            if btn.objectName() == "BTN_REMS":                                          ## Lock
                btn.setObjectName("BTN_REMS_LOCKED")
                btn.setStyleSheet(QSS.getStylesheet())
                btn.setToolTip(self.TTIP_BTN_REMS_LOCKED)
                self.CBX_RLS[i].setEnabled(False)
                self.CBX_NMS[i].setEnabled(False)
            else:                                                                       ## Unlock
                btn.setObjectName("BTN_REMS")
                btn.setStyleSheet(QSS.getStylesheet())
                btn.setToolTip(self.TTIP_BTN_REMS)
                self.CBX_RLS[i].setEnabled(True)
                self.CBX_NMS[i].setEnabled(True)


    def mouseReleased(self, button, event):
        """
        Overrides default mouse release events for BTN_INSS and BTN_REMS
        Handles mouse-related events from field buttons such as:
        - Insert
        - Remove
        """

        ## Insert Button
        if button == 'INSS':
            if event.button() == Qt.LeftButton:
                self.redirectFieldInsertion()
            elif event.button() == Qt.RightButton:
                self.redirectFieldInsertion(True)

        ## Remove Button
        elif button == 'REMS':
            if event.button() == Qt.LeftButton:
                self.removeField()
            elif event.button() == Qt.RightButton:
                self.lockUnlockField()


    def ignoreWheel(self, event):
        pass
    

    def recordCbx(self, cbx):
        """
        Overrides TextChanged event of combo boxes 

        Records the CBX object to help identifying the last
        combo box user used for editing fields
        
        PREV_CBX = (Category, Index Pos, Current Index Pos, Current Text)
        """
        if self.LOADING: return

        if cbx == 'RLS':
            for i, cbx in enumerate(self.CBX_RLS):
                if not cbx.hasFocus(): continue
                if cbx.currentText() == self.SEPARATOR:                                                             ## Helps preventing to display the separator
                    cbx.setCurrentIndex(cbx.currentIndex()+ (1 if self.PREV_CBX[2] < cbx.currentIndex() else -1))
                    return
                self.PREV_CBX = (0, i, cbx.currentIndex(), cbx.itemText(cbx.currentIndex()))
                break
                
        elif cbx == 'NMS':
            for i, cbx in enumerate(self.CBX_NMS):
                if not cbx.hasFocus(): continue
                if cbx.currentText() == self.SEPARATOR:
                    cbx.setCurrentIndex(cbx.currentIndex()+ (1 if self.PREV_CBX[2] < cbx.currentIndex() else -1))
                    return
                self.PREV_CBX = (1, i, cbx.currentIndex(), cbx.itemText(cbx.currentIndex()))
                break


    def getFieldData(self, merge=False):
        """
        Returns 2 lists, and a dictionary of all fields' data
        Used for when exporting to a file (Powerpoint, plain text)
        """
        ROLES = [role.currentText() for role in self.CBX_RLS]
        NAMES = [name.currentText() for name in self.CBX_NMS]
        DICT = {i:[k,v] for i, (k,v) in enumerate(zip(ROLES, NAMES))}

        if merge:
            ROLES += [i for i in RLS if i not in ROLES]
            NAMES += [i for i in NMS if i not in NAMES]
        
        return ROLES, NAMES, DICT

    
    def checkKeyPress(self, e):
        """
        Handles forwarded KeyPressEvent from UIA for field switching via Enter|Return key
        """
        TGT = self.CBX_NMS if self.PREV_CBX[0] else self.CBX_RLS
        for i, cbx in enumerate(TGT):
            if not cbx.hasFocus(): continue
            TGT[i+1 if i+1 != len(TGT) else 0].setFocus()
            break




class Export(object):
    """
    Handles exporting related functions.
    """
    def __init__(self):
        pass

    
    def fromActiveField(self, i, clear=False):
        """
        Exports file from an active field
        """
        if clear: R, N = '', ''
        else: R, N = FLD.CBX_RLS[i].currentText(), FLD.CBX_NMS[i].currentText() 

        try:
            with open(f"{SYS.DIR_PROGRAM}/Role.txt", "w") as f: f.write(R)
            with open(f"{SYS.DIR_PROGRAM}/Name.txt", "w") as f: f.write(N)
            # if not clear: LOG.info(f"Active field exported: {R} - {N}")
        except Exception as e:
            LOG.error(e)
            


    def toPlainText(self):
        """
        Export current values into a plain text
        """
        RLS, NMS, DCT = FLD.getFieldData()
        
        DIR_TGT = QtWidgets.QFileDialog.getExistingDirectory(None, "Select Folder", PKG.DIR_EXPORT_RECENT)
        if DIR_TGT:
            LOG.info(f"Saving file to {DIR_TGT}")
            try:
                with open(f"{DIR_TGT}/Roles.txt", "w") as f: f.write('\n'.join(RLS))
                with open(f"{DIR_TGT}/Names.txt", "w") as f: f.write('\n'.join(NMS))
                with open(f"{DIR_TGT}/Participants.txt", "w") as f: f.write('\n'.join([f"{r}: {n}" for r,n in zip(RLS,NMS)]))
            except Exception as e:
                LOG.error(e)
            
            ## Save to JSON
            RLS, NMS, DCT = FLD.getFieldData(True)
            DCFG['POOL'].update({"ROLES": RLS})
            DCFG['POOL'].update({"NAMES": NMS})
            DCFG['CONFIG'].update({"DIR_EXPORT_RECENT": DIR_TGT})
            PKG.DIR_EXPORT_RECENT = DIR_TGT
            PDB.dump()
            LOG.info("File successfully saved.")
            ## 
            ## Add code here that would update the combo boxes connected to FDS.fillupItems()
            ## 


    def toPowerpoint(self):
        """
        Export data to a Powerpoint file
        """
        def singleColumn(self):
            RLS, NMS, DCT = FLD.getFieldData()
            self.PRS = Presentation()

            ## 16:9 Ratio
            self.PRS.slide_width, self.PRS.slide_height = Inches(16), Inches(9)

            SLD_MAIN = self.PRS.slides.add_slide(self.PRS.slide_layouts[6])
            SLD_MAIN.shapes.add_picture(PKG.IMG_BACKGROUND, 0, 0, self.PRS.slide_width, self.PRS.slide_height) # Left-Top-Width-Height

            ## Textboxes
            ## Title
            TBX_TITLE = SLD_MAIN.shapes.add_textbox(0, Cm(1), self.PRS.slide_width, Cm(4))
            FRM_TITLE = TBX_TITLE.text_frame

            PARA_TITLE = FRM_TITLE.paragraphs[0]
            PARA_TITLE.text = PKG.TXT_TITLE
            PARA_TITLE.font.size = PKG.FSZ_TITLE
            PARA_TITLE.font.color.rgb = PKG.RGB_TITLE
            PARA_TITLE.alignment = PP_ALIGN.CENTER
            PARA_TITLE.font.name = PKG.FONT_TITLE
            PARA_TITLE.font.italic = True

            ## Subtitle
            TBX_SUBTITLE = SLD_MAIN.shapes.add_textbox(0, Cm(3), self.PRS.slide_width, Cm(4))
            FRM_SUBTITLE = TBX_SUBTITLE.text_frame
            PARA_SUBTITLE = FRM_SUBTITLE.paragraphs[0]
            PARA_SUBTITLE.text = PKG.TXT_SUBTITLE
            PARA_SUBTITLE.font.size = PKG.FSZ_SUBTITLE
            PARA_SUBTITLE.font.color.rgb = PKG.RGB_SUBTITLE
            PARA_SUBTITLE.alignment = PP_ALIGN.CENTER
            PARA_SUBTITLE.font.name = PKG.FONT_SUBTITLE

            ## Date
            TBX_DATE = SLD_MAIN.shapes.add_textbox(0, Cm(4.2), self.PRS.slide_width, Cm(4))
            FRM_DATE = TBX_DATE.text_frame
            PARA_DATE = FRM_DATE.paragraphs[0]
            PARA_DATE.text = f"——— {datetime.datetime.now().strftime('%B %d, %Y')} ———"
            PARA_DATE.font.size = PKG.FSZ_DATE
            PARA_DATE.font.color.rgb = PKG.RGB_DATE
            PARA_DATE.alignment = PP_ALIGN.CENTER
            PARA_DATE.font.name = PKG.FONT_DATE


            ## Paragraph
            left, top, width, height = 0, Cm(5.4), int(self.PRS.slide_width / 2), int(self.PRS.slide_height) - Cm(2)
            PARA_FNTSZ = Pt(39-FLD.FIELDS)

            ## Roles
            TBX_ROLES = SLD_MAIN.shapes.add_textbox(left, top, width, height)
            FRM_ROLES = TBX_ROLES.text_frame

            PARA_RLS = FRM_ROLES.paragraphs[0]
            PARA_RLS.text = '\n'.join([f"{r}:" for r in RLS])
            PARA_RLS.font.size = PARA_FNTSZ
            PARA_RLS.font.color.rgb = PKG.RGB_ROLES
            PARA_RLS.alignment = PP_ALIGN.RIGHT
            PARA_RLS.font.name = PKG.FONT_CONTENT
            PARA_RLS.font.italic = True

            ## Names
            TBX_NAMES = SLD_MAIN.shapes.add_textbox(int(self.PRS.slide_width / 2), top, width, height)
            FRM_NAMES = TBX_NAMES.text_frame

            PARA_NMS = FRM_NAMES.paragraphs[0]
            PARA_NMS.text = '\n'.join(NMS)
            PARA_NMS.font.size = PARA_FNTSZ
            PARA_NMS.font.color.rgb = PKG.RGB_NAMES
            PARA_NMS.alignment = PP_ALIGN.LEFT
            PARA_NMS.font.name = PKG.FONT_CONTENT


        def doubleColumns(self):
            RLS, NMS, DCT = FLD.getFieldData()
            self.PRS = Presentation()

            ## 16:9 Ratio
            self.PRS.slide_width, self.PRS.slide_height = Inches(16), Inches(9)

            SLD_MAIN = self.PRS.slides.add_slide(self.PRS.slide_layouts[6])
            SLD_MAIN.shapes.add_picture(PKG.IMG_BACKGROUND, 0, 0, self.PRS.slide_width, self.PRS.slide_height) # Left-Top-Width-Height

            ## Textboxes
            ## Title
            TBX_TITLE = SLD_MAIN.shapes.add_textbox(0, Cm(1), self.PRS.slide_width, Cm(4))
            FRM_TITLE = TBX_TITLE.text_frame

            PARA_TITLE = FRM_TITLE.paragraphs[0]
            PARA_TITLE.text = PKG.TXT_TITLE
            PARA_TITLE.font.size = PKG.FSZ_TITLE
            PARA_TITLE.font.color.rgb = PKG.RGB_TITLE
            PARA_TITLE.alignment = PP_ALIGN.CENTER
            PARA_TITLE.font.name = PKG.FONT_TITLE
            PARA_TITLE.font.italic = True

            ## Subtitle
            TBX_SUBTITLE = SLD_MAIN.shapes.add_textbox(0, Cm(3), self.PRS.slide_width, Cm(4))
            FRM_SUBTITLE = TBX_SUBTITLE.text_frame
            PARA_SUBTITLE = FRM_SUBTITLE.paragraphs[0]
            PARA_SUBTITLE.text = PKG.TXT_SUBTITLE
            PARA_SUBTITLE.font.size = PKG.FSZ_SUBTITLE
            PARA_SUBTITLE.font.color.rgb = PKG.RGB_SUBTITLE
            PARA_SUBTITLE.alignment = PP_ALIGN.CENTER
            PARA_SUBTITLE.font.name = PKG.FONT_SUBTITLE

            ## Sabbath School Text
            TBX_SUBTITLE = SLD_MAIN.shapes.add_textbox(Inches(4), Inches(3), self.PRS.slide_width, Cm(4))
            FRM_SUBTITLE = TBX_SUBTITLE.text_frame
            PARA_SUBTITLE = FRM_SUBTITLE.paragraphs[0]
            PARA_SUBTITLE.text = "Sabbath School"
            PARA_SUBTITLE.font.size = Pt(40)
            PARA_SUBTITLE.font.color.rgb = PKG.RGB_SUBTITLE
            PARA_SUBTITLE.alignment = PP_ALIGN.CENTER
            PARA_SUBTITLE.font.name = PKG.FONT_SUBTITLE

            ## Divine Service Text
            TBX_SUBTITLE = SLD_MAIN.shapes.add_textbox(Inches(1.25), Inches(6.5), Inches(4), Inches(2))
            FRM_SUBTITLE = TBX_SUBTITLE.text_frame
            PARA_SUBTITLE = FRM_SUBTITLE.paragraphs[0]
            PARA_SUBTITLE.text = "Divine Service"
            PARA_SUBTITLE.font.size = Pt(40)
            PARA_SUBTITLE.font.color.rgb = PKG.RGB_SUBTITLE
            PARA_SUBTITLE.alignment = PP_ALIGN.CENTER
            PARA_SUBTITLE.font.name = PKG.FONT_SUBTITLE


            ## Date
            TBX_DATE = SLD_MAIN.shapes.add_textbox(0, Cm(4.2), self.PRS.slide_width, Cm(4))
            FRM_DATE = TBX_DATE.text_frame
            PARA_DATE = FRM_DATE.paragraphs[0]
            PARA_DATE.text = f"——— {datetime.datetime.now().strftime('%B %d, %Y')} ———"
            PARA_DATE.font.size = PKG.FSZ_DATE
            PARA_DATE.font.color.rgb = PKG.RGB_DATE
            PARA_DATE.alignment = PP_ALIGN.CENTER
            PARA_DATE.font.name = PKG.FONT_DATE


            ## Paragraph
            left, top, width, height = 0, Cm(5.4), int(self.PRS.slide_width/2), int(self.PRS.slide_height) - Cm(2)
            PARA_FNTSZ = Pt(38-FLD.FIELDS)

            ## Sabbath School
            ## Roles
            TBX_ROLES = SLD_MAIN.shapes.add_textbox(0, top, width/2.2, height)
            TBX_NAMES = SLD_MAIN.shapes.add_textbox(width/2.2, top, width/2.2, height)
            FRM_ROLES = TBX_ROLES.text_frame
            FRM_NAMES = TBX_NAMES.text_frame
            TXT_ROLES = '\n'.join([f"{r}:" for r in RLS])
            TXT_NAMES = '\n'.join(NMS)


            ## Service-based Splitter
            SPLITTED = CORE.splitContents(TXT_ROLES, TXT_NAMES)

            PARA_RLS = FRM_ROLES.paragraphs[0]
            PARA_RLS.text = '\n'.join(SPLITTED[0])
            PARA_RLS.font.size = PARA_FNTSZ
            PARA_RLS.font.color.rgb = PKG.RGB_ROLES
            PARA_RLS.alignment = PP_ALIGN.RIGHT
            PARA_RLS.font.name = PKG.FONT_CONTENT
            PARA_RLS.font.italic = True

            ## Names
            PARA_NMS = FRM_NAMES.paragraphs[0]
            PARA_NMS.text = '\n'.join(SPLITTED[1])
            PARA_NMS.font.size = PARA_FNTSZ
            PARA_NMS.font.color.rgb = PKG.RGB_NAMES
            PARA_NMS.alignment = PP_ALIGN.LEFT
            PARA_NMS.font.name = PKG.FONT_CONTENT
            

            ## Divine Worship
            ## Roles
            TBX_ROLES = SLD_MAIN.shapes.add_textbox(Inches(7.25), Inches(5.5), width/2.2, height)
            TBX_NAMES = SLD_MAIN.shapes.add_textbox(Inches(10.9), Inches(5.5), width/2.2, height)
            FRM_ROLES = TBX_ROLES.text_frame
            FRM_NAMES = TBX_NAMES.text_frame

            PARA_RLS = FRM_ROLES.paragraphs[0]
            PARA_RLS.text = '\n'.join(SPLITTED[2])
            PARA_RLS.font.size = PARA_FNTSZ
            PARA_RLS.font.color.rgb = PKG.RGB_ROLES
            PARA_RLS.alignment = PP_ALIGN.RIGHT
            PARA_RLS.font.name = PKG.FONT_CONTENT
            PARA_RLS.font.italic = True

            ## Names
            PARA_NMS = FRM_NAMES.paragraphs[0]
            PARA_NMS.text = '\n'.join(SPLITTED[3])
            PARA_NMS.font.size = PARA_FNTSZ
            PARA_NMS.font.color.rgb = PKG.RGB_NAMES
            PARA_NMS.alignment = PP_ALIGN.LEFT
            PARA_NMS.font.name = PKG.FONT_CONTENT
        

        ## Scan for new items that are currently not in pool
        ##
        ## NOTE: Needed to change the structure of data.json.
        ##       the NAMES pool will be extended into 2 (USED & XTRA)
        ##       USED - Names that are displayed in the fields. Same length with ROLES
        ##       XTRA - Unused names that are in the pool

        # NEW_NAMES = [name for name in [c.currentText() for c in FLD.CBX_NMS] if name not in DCFG['POOL']['NAMES']]
        # if len(NEW_NAMES):
        #     MSG_BOX = QtWidgets.QMessageBox()
        #     MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON))
        #     MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No)
        #     MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.No)
        #     MSG_BOX.setIcon(QtWidgets.QMessageBox.Question)
        #     MSG_BOX.setText("There are new participants recorded\nDo you want to add them all to your list?")
        #     MSG_BOX.setDetailedText("\n".join(map(str, NEW_NAMES)))
        #     MSG_BOX.setWindowTitle(SW.NAME)
        #     MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
        #     MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')
        #     RET = MSG_BOX.exec_()

        #     if RET == QtWidgets.QMessageBox.Yes:
        #         p('Saved')
        #     else:
        #         p('Unsaved')
        #         return

        START = time.time()
        LOG.info("Generating Powerpoint")
        doubleColumns(self) if FLD.FIELDS > FLD.FIELDS_MAX/1.15 else singleColumn(self)            

        try:
            self.PRS.save(SYS.FILE_PPT_EXPORTED)
            os.startfile(SYS.FILE_PPT_EXPORTED)
        except PermissionError as e:
            LOG.warn("PermissionError: The file is still open or is already running. Close the file first and try again.")
            MSG_BOX = QtWidgets.QMessageBox()
            MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON))
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Ok)
            MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.Ok)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Warning)
            MSG_BOX.setText("The file is still open or is already running.\nIf you have changes, close the file first and try again.")
            MSG_BOX.setWindowTitle(SW.NAME)
            MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
            MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')
            MSG_BOX.exec_()
        except Exception as e:
            LOG.error(f"{e}")
        else:
            RLS, NMS, DCT = FLD.getFieldData(True)
            DCFG['POOL'].update({"ROLES": RLS})
            DCFG['POOL'].update({"NAMES": NMS})
            PDB.dump()
            LOG.info(f"File successfully saved. ({round((time.time()-START)*1000)} ms)")
            UIB.hide()




class QWGT_PARTICIPANTS(QtWidgets.QMainWindow):
    """
    Main user interface window for this software.
    Uses PyQt5 for handling GUI.

    Alias UIA (User Interface A / Main)
    """
    def __init__(self, parent = None):
        QtWidgets.QMainWindow.__init__(self, parent)
    

    def setupUI(self):
        """
        Initializes all widgets for UIA
        """
        ## Window
        self.setObjectName("WIN_PARTICIPANTS")
        self.resize(0, 0)
        self.setMaximumWidth(520)
        self.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON))
        self.setWindowFlags(self.windowFlags() | QtCore.Qt.WindowStaysOnTopHint)

        ## Layouts
        self.WGT_CENTRAL = QtWidgets.QWidget(self); self.WGT_CENTRAL.setObjectName("WGT_CENTRAL")
        self.LYT_MAIN = QtWidgets.QGridLayout(self.WGT_CENTRAL); self.LYT_MAIN.setObjectName("LYT_MAIN")
        self.LYT_HEAD = QtWidgets.QGridLayout(); self.LYT_HEAD.setObjectName("LYT_HEAD")
        self.LYT_BODY = QtWidgets.QGridLayout(); self.LYT_BODY.setObjectName("LYT_BODY")
        self.LYT_FOOTER = QtWidgets.QGridLayout(); self.LYT_FOOTER.setObjectName("LYT_FOOTER")

        self.LYT_ROLES = QtWidgets.QVBoxLayout(); self.LYT_BODY.setObjectName("LYT_ROLES")
        self.LYT_NAMES = QtWidgets.QVBoxLayout(); self.LYT_BODY.setObjectName("LYT_NAMES")
        self.LYT_CLEAR = QtWidgets.QVBoxLayout(); self.LYT_BODY.setObjectName("LYT_CLEAR")
        self.LYT_INSRT = QtWidgets.QVBoxLayout(); self.LYT_BODY.setObjectName("LYT_INSRT")
        self.LYT_ACTIV = QtWidgets.QVBoxLayout(); self.LYT_BODY.setObjectName("LYT_ACTIV")
        
        ## Banner
        self.PIX_HEADER = QtWidgets.QLabel(self.WGT_CENTRAL); self.PIX_HEADER.setObjectName("PIX_HEADER"); self.PIX_HEADER.setAlignment(QtCore.Qt.AlignCenter)        
        self.PIX_HEADER.setPixmap(QtGui.QPixmap(SYS.RES_HEADERLOGO).scaledToHeight(75, QtCore.Qt.KeepAspectRatio & Qt.SmoothTransformation))

        ## Head
        SPC_WINV_TOPP = QtWidgets.QSpacerItem(20, 15, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum)
        SPC_TITLEV2 = QtWidgets.QSpacerItem(20, 15, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum)
        self.BTN_SETTINGS = QtWidgets.QPushButton(self.WGT_CENTRAL); self.BTN_SETTINGS.setObjectName("BTN_SETTINGS")
        self.BTN_SAVELIST = QtWidgets.QPushButton(self.WGT_CENTRAL); self.BTN_SAVELIST.setObjectName("BTN_SAVELIST")
        SPC_HEADH1 = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)
        self.LBL_EXPORT = QtWidgets.QLabel(self.WGT_CENTRAL); self.LBL_EXPORT.setObjectName("LBL_EXPORT")
        self.BTN_POWERPOINT = QtWidgets.QPushButton(self.WGT_CENTRAL); self.BTN_POWERPOINT.setObjectName("BTN_POWERPOINT")
        self.BTN_PLAINTEXT = QtWidgets.QPushButton(self.WGT_CENTRAL); self.BTN_PLAINTEXT.setObjectName("BTN_PLAINTEXT")
        
        ## Body
        SPC_BODYV1 = QtWidgets.QSpacerItem(20, 15, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum)
        self.LBL_RL = QtWidgets.QLabel(self.WGT_CENTRAL); self.LBL_RL.setObjectName("LBL_RL"); self.LBL_RL.setAlignment(QtCore.Qt.AlignCenter)
        SPC_BODYH = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)
        self.LBL_NM = QtWidgets.QLabel(self.WGT_CENTRAL); self.LBL_NM.setObjectName("LBL_NM"); self.LBL_NM.setAlignment(QtCore.Qt.AlignCenter)
        SPC_BODYV2 = QtWidgets.QSpacerItem(20, 20, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Minimum)
        
        ## Window Spacers
        SPC_WINH_LEFT = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)
        SPC_WINV_BOTM = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_WINH_RGHT = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)

        ## Layering
        ## Head
        self.LYT_HEAD.addWidget(self.PIX_HEADER, 1, 0, 1, 8)
        self.LYT_HEAD.addItem(SPC_TITLEV2, 2, 0, 1, 3)
        self.LYT_HEAD.addWidget(self.BTN_SETTINGS, 3, 0, 1, 1)
        self.LYT_HEAD.addWidget(self.BTN_SAVELIST, 3, 1, 1, 1)
        self.LYT_HEAD.addItem(SPC_HEADH1, 3, 3, 1, 1)
        self.LYT_HEAD.addWidget(self.LBL_EXPORT, 3, 4, 1, 1)
        self.LYT_HEAD.addWidget(self.BTN_POWERPOINT, 3, 5, 1, 1)
        self.LYT_HEAD.addWidget(self.BTN_PLAINTEXT, 3, 6, 1, 1)
        self.LYT_HEAD.addItem(SPC_BODYV1, 4, 0, 1, 3)
        self.LYT_HEAD.addLayout(self.LYT_BODY, 5, 0, 1, 7)
        self.LYT_HEAD.addItem(SPC_BODYV2, 6, 0, 1, 3)

        ## Body
        self.LYT_BODY.addWidget(self.LBL_RL, 0, 0, 1, 1)
        self.LYT_BODY.addLayout(self.LYT_ROLES, 1, 0, 1, 1)
        self.LYT_BODY.addItem(SPC_BODYH, 0, 1, 1, 1)
        self.LYT_BODY.addWidget(self.LBL_NM, 0, 2, 1, 1)
        self.LYT_BODY.addLayout(self.LYT_NAMES, 1, 2, 1, 1)
        self.LYT_BODY.addLayout(self.LYT_ACTIV, 1, 3, 1, 1)
        self.LYT_BODY.addLayout(self.LYT_INSRT, 1, 4, 1, 1)
        self.LYT_BODY.addLayout(self.LYT_CLEAR, 1, 5, 1, 1)

        ## Window
        self.LYT_MAIN.addLayout(self.LYT_HEAD, 0, 1, 1, 1)
        self.LYT_HEAD.addItem(SPC_WINV_TOPP, 0, 0, 1, 3)
        self.LYT_MAIN.addItem(SPC_WINH_LEFT, 0, 0, 1, 1)
        self.LYT_MAIN.addItem(SPC_WINV_BOTM, 1, 1, 1, 1)
        self.LYT_MAIN.addItem(SPC_WINH_RGHT, 0, 2, 1, 1)
        self.setCentralWidget(self.WGT_CENTRAL)

        ## Stylesheets
        QSS.setupFonts()

        ## Initialization
        self.setupDisplay()
        self.setupConnections()
        FLD.setup()
    

    def setupDisplay(self):
        """
        Handle all UI-based properties including texts
        Identical to `retranslateUi` method.
        """
        ## Main
        self.setWindowTitle("Participants")

        ## Head
        self.PIX_HEADER.setToolTip(f"v{SW.VERSION} {SW.VERSION_NAME}\nMade for Seventh-day Adventist Church\n\n© {SW.PROD_YEAR} {SW.AUTHOR}")
        self.LBL_EXPORT.setText("Export to")

        ## Body
        self.LBL_RL.setText("ROLE")
        self.LBL_NM.setText("NAME")
    

    def setupConnections(self):
        """
        Manages Signals and Slots from user interactions
        """
        self.BTN_PLAINTEXT.clicked.connect(lambda: EXP.toPlainText())
        self.BTN_POWERPOINT.clicked.connect(lambda: EXP.toPowerpoint())
        self.BTN_SETTINGS.clicked.connect(lambda: UIB.enterWindow())
        

    def updateWindowTitle(self):
        """
        Handles UIA's main window title
        """
        self.setWindowTitle(f"{SW.NAME} ({FLD.FIELDS})")


    def keyPressEvent(self, e):
        if e.key() in (Qt.Key_Enter, Qt.Key_Return):
            FLD.checkKeyPress(e)
    
    



    
class QWGT_SETTINGS(QtWidgets.QMainWindow):
    """
    Settings window 
    Mostly for UI-related functions only.

    Alias UIB (User Interface B / Settings)
    """
    def __init__(self, parent = None):
        QtWidgets.QWidget.__init__(self, parent)
        self.CHANGED = False
        self.BG_SCALING = 100


    def setupUI(self):
        """
        Initializes all widgets for UIB
        """
        ## Window
        self.setObjectName("WIN_SETTINGS")
        self.setFixedSize(380, 720)
        self.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON))
        self.setWindowFlags(self.windowFlags() | QtCore.Qt.WindowStaysOnTopHint)

        self.WGT_CENTRAL = QtWidgets.QWidget(self); self.WGT_CENTRAL.setObjectName("WGT_CENTRAL")
        self.setCentralWidget(self.WGT_CENTRAL)
        
        ## Footer
        self.LBL_APPVERSION = QtWidgets.QLabel(self.WGT_CENTRAL)
        self.LBL_APPVERSION.setEnabled(False); self.LBL_APPVERSION.setObjectName("LBL_APPVERSION")
        self.BTN_OK = QtWidgets.QPushButton(self.WGT_CENTRAL); self.BTN_OK.setObjectName("BTN_OK")

        ## Settings Main Tab
        self.TBW_SETTINGS = QtWidgets.QTabWidget(self.WGT_CENTRAL); self.TBW_SETTINGS.setObjectName("TBW_SETTINGS")

        ## General Tab
        self.TAB_GENERAL = QtWidgets.QWidget(); self.TAB_GENERAL.setObjectName("TAB_GENERAL")
        self.GBX_GEN_APPEARANCE = QtWidgets.QGroupBox(self.TAB_GENERAL); self.GBX_GEN_APPEARANCE.setObjectName("GBX_GEN_APPEARANCE")
        self.CHK_GEN_ALWAYS_ON_TOP = QtWidgets.QCheckBox(self.GBX_GEN_APPEARANCE); self.CHK_GEN_ALWAYS_ON_TOP.setObjectName("CHK_GEN_ALWAYS_ON_TOP")
        self.CHK_GEN_USE_SUGGESTED = QtWidgets.QCheckBox(self.GBX_GEN_APPEARANCE); self.CHK_GEN_USE_SUGGESTED.setObjectName("CHK_GEN_USE_SUGGESTED")

        self.GBX_GEN_PRESETS = QtWidgets.QGroupBox(self.TAB_GENERAL); self.GBX_GEN_PRESETS.setObjectName("GBX_GEN_PRESETS")
        self.LST_GEN_PRESETS = QtWidgets.QListWidget(self.GBX_GEN_PRESETS); self.LST_GEN_PRESETS.setObjectName("LST_GEN_PRESETS")
        self.BTN_GEN_MODIFY = QtWidgets.QPushButton(self.GBX_GEN_PRESETS); self.BTN_GEN_MODIFY.setObjectName("BTN_GEN_MODIFY")
        self.BTN_GEN_IMPORT = QtWidgets.QPushButton(self.GBX_GEN_PRESETS); self.BTN_GEN_IMPORT.setObjectName("BTN_GEN_IMPORT")
        self.BTN_GEN_REMOVE = QtWidgets.QPushButton(self.GBX_GEN_PRESETS); self.BTN_GEN_REMOVE.setObjectName("BTN_GEN_REMOVE")

        ## Presentation Tab
        self.TAB_PRESENTATION = QtWidgets.QWidget(); self.TAB_PRESENTATION.setObjectName("TAB_PRESENTATION")

        ## Presentation - Content
        self.GBX_PRES_CONTENT = QtWidgets.QGroupBox(self.TAB_PRESENTATION); self.GBX_PRES_CONTENT.setObjectName("GBX_PRES_CONTENT")
        self.LBL_CNTT_TITLE = QtWidgets.QLabel(self.GBX_PRES_CONTENT); self.LBL_CNTT_TITLE.setObjectName("LBL_CNTT_TITLE"); self.LBL_CNTT_TITLE.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.LNE_CNTT_TITLE = QtWidgets.QLineEdit(self.GBX_PRES_CONTENT); self.LNE_CNTT_TITLE.setObjectName("LNE_CNTT_TITLE"); self.LNE_CNTT_TITLE.setClearButtonEnabled(True)
        self.LBL_CNTT_SUBTITLE = QtWidgets.QLabel(self.GBX_PRES_CONTENT); self.LBL_CNTT_SUBTITLE.setObjectName("LBL_CNTT_SUBTITLE"); self.LBL_CNTT_SUBTITLE.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.LNE_CNTT_SUBTITLE = QtWidgets.QLineEdit(self.GBX_PRES_CONTENT); self.LNE_CNTT_SUBTITLE.setObjectName("LNE_CNTT_SUBTITLE"); self.LNE_CNTT_SUBTITLE.setClearButtonEnabled(True)
        self.CHK_CNTT_DISPDATE = QtWidgets.QCheckBox(self.GBX_PRES_CONTENT); self.CHK_CNTT_DISPDATE.setObjectName("CHK_CNTT_DISPDATE")
        self.CHK_CNTT_USEWIDESCR = QtWidgets.QCheckBox(self.GBX_PRES_CONTENT); self.CHK_CNTT_USEWIDESCR.setObjectName("CHK_CNTT_USEWIDESCR")
        self.CHK_CNTT_ENABLECBXSCROLL = QtWidgets.QCheckBox(self.GBX_PRES_CONTENT); self.CHK_CNTT_ENABLECBXSCROLL.setObjectName("CHK_CNTT_ENABLECBXSCROLL")
        
        ## Presentation - Fonts and Colors
        self.GBX_PRES_FONTSCOLORS = QtWidgets.QGroupBox(self.TAB_PRESENTATION); self.GBX_PRES_FONTSCOLORS.setObjectName("GBX_PRES_FONTSCOLORS")
        self.LBL_FAC_TITLE = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_TITLE.setObjectName("LBL_FAC_TITLE"); self.LBL_FAC_TITLE.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.LBL_FAC_SUBTITLE = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_SUBTITLE.setObjectName("LBL_FAC_SUBTITLE"); self.LBL_FAC_SUBTITLE.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.LBL_FAC_BODY = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_BODY.setObjectName("LBL_FAC_BODY"); self.LBL_FAC_BODY.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.LBL_FAC_ROLE = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_ROLE.setObjectName("LBL_FAC_ROLE"); self.LBL_FAC_ROLE.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.LBL_FAC_NAME = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_NAME.setObjectName("LBL_FAC_NAME"); self.LBL_FAC_NAME.setAlignment(QtCore.Qt.AlignRight|QtCore.Qt.AlignTrailing|QtCore.Qt.AlignVCenter)
        self.FCB_FAC_TITLE = QtWidgets.QFontComboBox(self.GBX_PRES_FONTSCOLORS); self.FCB_FAC_TITLE.setObjectName("FCB_FAC_TITLE"); self.FCB_FAC_TITLE.setMaximumWidth(140)
        self.FCB_FAC_SUBITITLE = QtWidgets.QFontComboBox(self.GBX_PRES_FONTSCOLORS); self.FCB_FAC_SUBITITLE.setObjectName("FCB_FAC_SUBITITLE"); self.FCB_FAC_SUBITITLE.setMaximumWidth(140)
        self.FCB_FAC_BODY = QtWidgets.QFontComboBox(self.GBX_PRES_FONTSCOLORS); self.FCB_FAC_BODY.setObjectName("FCB_FAC_BODY"); self.FCB_FAC_BODY.setMaximumWidth(140)
        self.FCB_FAC_ROLE = QtWidgets.QFontComboBox(self.GBX_PRES_FONTSCOLORS); self.FCB_FAC_ROLE.setObjectName("FCB_FAC_ROLE"); self.FCB_FAC_ROLE.setMaximumWidth(140)
        self.FCB_FAC_NAME = QtWidgets.QFontComboBox(self.GBX_PRES_FONTSCOLORS); self.FCB_FAC_NAME.setObjectName("FCB_FAC_NAME"); self.FCB_FAC_NAME.setMaximumSize(QtCore.QSize(140, 16777215))
        self.LIN_FAC_DIVIDER = QtWidgets.QFrame(self.GBX_PRES_FONTSCOLORS); self.LIN_FAC_DIVIDER.setObjectName("LIN_FAC_DIVIDER"); self.LIN_FAC_DIVIDER.setFrameShape(QtWidgets.QFrame.HLine); self.LIN_FAC_DIVIDER.setFrameShadow(QtWidgets.QFrame.Sunken)
        self.LBL_FAC_TITLE_PREV = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_TITLE_PREV.setObjectName("LBL_FAC_TITLE_PREV")
        self.LBL_FAC_SUBTITLE_PREV = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_SUBTITLE_PREV.setObjectName("LBL_FAC_SUBTITLE_PREV")
        self.LBL_FAC_BODY_PREV = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_BODY_PREV.setObjectName("LBL_FAC_BODY_PREV")
        self.LBL_FAC_ROLE_PREV = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_ROLE_PREV.setObjectName("LBL_FAC_ROLE_PREV")
        self.LBL_FAC_NAME_PREV = QtWidgets.QLabel(self.GBX_PRES_FONTSCOLORS); self.LBL_FAC_NAME_PREV.setObjectName("LBL_FAC_NAME_PREV")
        self.BTN_FAC_TITLE = QtWidgets.QPushButton(self.GBX_PRES_FONTSCOLORS); self.BTN_FAC_TITLE.setObjectName("BTN_FAC_COLORPICKER"); self.BTN_FAC_TITLE.setMaximumSize(QtCore.QSize(30, 16777215))
        self.BTN_FAC_SUBTITLE = QtWidgets.QPushButton(self.GBX_PRES_FONTSCOLORS); self.BTN_FAC_SUBTITLE.setObjectName("BTN_FAC_COLORPICKER"); self.BTN_FAC_SUBTITLE.setMaximumSize(QtCore.QSize(30, 16777215))
        self.BTN_FAC_BODY = QtWidgets.QPushButton(self.GBX_PRES_FONTSCOLORS); self.BTN_FAC_BODY.setObjectName("BTN_FAC_COLORPICKER"); self.BTN_FAC_BODY.setMaximumSize(QtCore.QSize(30, 16777215))
        self.BTN_FAC_ROLE = QtWidgets.QPushButton(self.GBX_PRES_FONTSCOLORS); self.BTN_FAC_ROLE.setObjectName("BTN_FAC_COLORPICKER"); self.BTN_FAC_ROLE.setMaximumSize(QtCore.QSize(30, 16777215))
        self.BTN_FAC_NAME = QtWidgets.QPushButton(self.GBX_PRES_FONTSCOLORS); self.BTN_FAC_NAME.setObjectName("BTN_FAC_COLORPICKER"); self.BTN_FAC_NAME.setMaximumSize(QtCore.QSize(30, 16777215))
        
        ## Presentation - Background
        self.GBX_PRES_BACKGROUND = QtWidgets.QGroupBox(self.TAB_PRESENTATION); self.GBX_PRES_BACKGROUND.setObjectName("GBX_PRES_BACKGROUND")
        self.BTN_BG_BROWSE = QtWidgets.QPushButton(self.GBX_PRES_BACKGROUND); self.BTN_BG_BROWSE.setObjectName("BTN_BG_BROWSE"); self.BTN_BG_BROWSE.setMaximumSize(QtCore.QSize(30, 16777215))
        self.BTN_BG_DISCARD = QtWidgets.QPushButton(self.GBX_PRES_BACKGROUND); self.BTN_BG_DISCARD.setObjectName("BTN_BG_DISCARD"); self.BTN_BG_DISCARD.setMaximumSize(QtCore.QSize(30, 16777215))
        self.LBL_BG_PREVIEW = QtWidgets.QLabel(self.GBX_PRES_BACKGROUND); self.LBL_BG_PREVIEW.setObjectName("LBL_BG_PREVIEW"); self.LBL_BG_PREVIEW.setAlignment(QtCore.Qt.AlignCenter)
        self.PIX_BG_PREVIEW = QtWidgets.QLabel(self.GBX_PRES_BACKGROUND); self.PIX_BG_PREVIEW.setObjectName("PIX_BG_PREVIEW"); self.PIX_BG_PREVIEW.setAlignment(QtCore.Qt.AlignCenter)        
        self.PIX_BG_PREVIEW.setPixmap(QtGui.QPixmap(PKG.IMG_BACKGROUND).scaledToHeight(self.BG_SCALING, QtCore.Qt.KeepAspectRatio & Qt.SmoothTransformation))
        
        self.BTN_PRES_RESET = QtWidgets.QPushButton(self.TAB_PRESENTATION); self.BTN_PRES_RESET.setObjectName("BTN_PRES_RESET")

        ## Members Tab
        self.TAB_MEMBERS = QtWidgets.QWidget(); self.TAB_MEMBERS.setObjectName("TAB_MEMBERS")
        self.LNE_MEM_SEARCHADD = QtWidgets.QLineEdit(self.TAB_MEMBERS); self.LNE_MEM_SEARCHADD.setObjectName("LNE_MEM_SEARCHADD"); self.LNE_MEM_SEARCHADD.setClearButtonEnabled(True)
        self.LST_MEM_MEMBERS = QtWidgets.QListWidget(self.TAB_MEMBERS); self.LST_MEM_MEMBERS.setObjectName("LST_MEM_MEMBERS")
        self.BTN_MEM_ADD = QtWidgets.QPushButton(self.TAB_MEMBERS); self.BTN_MEM_ADD.setObjectName("BTN_MEM_ADD")
        self.BTN_MEM_EDIT = QtWidgets.QPushButton(self.TAB_MEMBERS); self.BTN_MEM_EDIT.setObjectName("BTN_MEM_EDIT"); self.BTN_MEM_EDIT.setEnabled(False)
        self.BTN_MEM_REMOVE = QtWidgets.QPushButton(self.TAB_MEMBERS); self.BTN_MEM_REMOVE.setObjectName("BTN_MEM_REMOVE"); self.BTN_MEM_REMOVE.setEnabled(False)
        self.BTN_MEM_IMPORT = QtWidgets.QPushButton(self.TAB_MEMBERS); self.BTN_MEM_IMPORT.setObjectName("BTN_MEM_IMPORT")
        self.BTN_MEM_EXPORT = QtWidgets.QPushButton(self.TAB_MEMBERS); self.BTN_MEM_EXPORT.setObjectName("BTN_MEM_EXPORT")
        self.GBX_MEM_DETAILS = QtWidgets.QGroupBox(self.TAB_MEMBERS);  self.GBX_MEM_DETAILS.setObjectName("GBX_MEM_DETAILS")
        self.LBL_DET_MEMBERS = QtWidgets.QLabel(self.GBX_MEM_DETAILS); self.LBL_DET_MEMBERS.setObjectName("LBL_DET_MEMBERS")
        self.LBL_DET_MEN = QtWidgets.QLabel(self.GBX_MEM_DETAILS); self.LBL_DET_MEN.setObjectName("LBL_DET_MEN")
        self.LBL_DET_WOMEN = QtWidgets.QLabel(self.GBX_MEM_DETAILS); self.LBL_DET_WOMEN.setObjectName("LBL_DET_WOMEN")
        self.LBL_DET_OTHERS = QtWidgets.QLabel(self.GBX_MEM_DETAILS); self.LBL_DET_OTHERS.setObjectName("LBL_DET_OTHERS")
        self.BTN_MEM_SAVE = QtWidgets.QPushButton(self.TAB_MEMBERS); self.BTN_MEM_SAVE.setObjectName("BTN_MEM_SAVE")

        ## Spacers
        SPC_FOOTER_H = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)
        SPC_GENERAL_V = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_GEN_PRESETS_V = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_PRES_V4 = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)
        SPC_PRES_BG_H = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_PRES_FAC_H = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Expanding, QtWidgets.QSizePolicy.Minimum)
        SPC_PRES_CNTT_V = QtWidgets.QSpacerItem(40, 20, QtWidgets.QSizePolicy.Preferred, QtWidgets.QSizePolicy.Minimum)
        SPC_PRES_V3 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_PRES_V2 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_PRES_V1 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
        SPC_MEM_V1 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)

        ## Grid Layouts
        self.GRID_MEMBERS = QtWidgets.QGridLayout(self.TAB_MEMBERS); self.GRID_MEMBERS.setObjectName("GRID_MEMBERS")
        self.GRID_MAIN = QtWidgets.QGridLayout(self.WGT_CENTRAL); self.GRID_MAIN.setObjectName("GRID_MAIN")
        self.GRID_GENERAL = QtWidgets.QGridLayout(self.TAB_GENERAL); self.GRID_GENERAL.setObjectName("GRID_GENERAL")
        self.GRID_PRES_BG = QtWidgets.QGridLayout(self.GBX_PRES_BACKGROUND); self.GRID_PRES_BG.setObjectName("GRID_PRES_BG")
        self.GRID_PRESENTATION = QtWidgets.QGridLayout(self.TAB_PRESENTATION); self.GRID_PRESENTATION.setObjectName("GRID_PRESENTATION")
        self.GRID_PRES_FAC = QtWidgets.QGridLayout(self.GBX_PRES_FONTSCOLORS); self.GRID_PRES_FAC.setObjectName("GRID_PRES_FAC")
        self.GRID_GEN_APPEARANCE = QtWidgets.QGridLayout(self.GBX_GEN_APPEARANCE); self.GRID_GEN_APPEARANCE.setObjectName("GRID_GEN_APPEARANCE")
        self.GRID_PRES_CONTENT = QtWidgets.QGridLayout(self.GBX_PRES_CONTENT); self.GRID_PRES_CONTENT.setObjectName("GRID_PRES_CONTENT")
        self.GRID_MEM_DET = QtWidgets.QGridLayout(self.GBX_MEM_DETAILS); self.GRID_MEM_DET.setObjectName("GRID_MEM_DET")
        self.GRID_GEN_PRESETS = QtWidgets.QGridLayout(self.GBX_GEN_PRESETS); self.GRID_GEN_PRESETS.setObjectName("GRID_GEN_PRESETS")

        ## Layering
        ## Main
        self.GRID_MAIN.addWidget(self.TBW_SETTINGS, 0, 0, 1, 3)
        self.GRID_MAIN.addWidget(self.LBL_APPVERSION, 1, 0, 1, 1)
        self.GRID_MAIN.addItem(SPC_FOOTER_H, 1, 1, 1, 1)
        self.GRID_MAIN.addWidget(self.BTN_OK, 1, 2, 1, 1)

        self.TBW_SETTINGS.addTab(self.TAB_GENERAL, "")
        self.TBW_SETTINGS.addTab(self.TAB_PRESENTATION, "")
        self.TBW_SETTINGS.addTab(self.TAB_MEMBERS, "")

        ## General Tab
        self.GRID_GENERAL.addWidget(self.GBX_GEN_APPEARANCE, 0, 0, 1, 1)
        self.GRID_GENERAL.addWidget(self.GBX_GEN_PRESETS, 1, 0, 1, 1)
        self.GRID_GENERAL.addItem(SPC_GENERAL_V, 2, 0, 1, 1)

        self.GRID_GEN_APPEARANCE.addWidget(self.CHK_GEN_ALWAYS_ON_TOP, 0, 0, 1, 1)
        self.GRID_GEN_APPEARANCE.addWidget(self.CHK_GEN_USE_SUGGESTED, 1, 0, 1, 1)
        self.GRID_GEN_APPEARANCE.addWidget(self.CHK_CNTT_ENABLECBXSCROLL, 2, 0, 1, 1)

        self.GRID_GEN_PRESETS.addWidget(self.LST_GEN_PRESETS, 0, 0, 6, 1)
        self.GRID_GEN_PRESETS.addWidget(self.BTN_GEN_IMPORT, 0, 1, 1, 1)
        self.GRID_GEN_PRESETS.addWidget(self.BTN_GEN_MODIFY, 1, 1, 1, 1)
        self.GRID_GEN_PRESETS.addWidget(self.BTN_GEN_REMOVE, 2, 1, 1, 1)
        self.GRID_GEN_PRESETS.addItem(SPC_GEN_PRESETS_V, 3, 1, 1, 1)
        
        ## Presentation Tab
        self.GRID_PRESENTATION.addWidget(self.GBX_PRES_CONTENT, 0, 0, 2, 2)
        self.GRID_PRESENTATION.addItem(SPC_PRES_V1, 2, 0, 1, 2)
        self.GRID_PRESENTATION.addWidget(self.GBX_PRES_FONTSCOLORS, 3, 0, 1, 2)
        self.GRID_PRESENTATION.addItem(SPC_PRES_V2, 4, 0, 1, 2)
        self.GRID_PRESENTATION.addWidget(self.GBX_PRES_BACKGROUND, 5, 0, 1, 2)
        self.GRID_PRESENTATION.addItem(SPC_PRES_V3, 6, 0, 1, 2)
        self.GRID_PRESENTATION.addWidget(self.BTN_PRES_RESET, 7, 0, 1, 1)
        self.GRID_PRESENTATION.addItem(SPC_PRES_V4, 7, 1, 1, 1)

        self.GRID_PRES_CONTENT.addWidget(self.LBL_CNTT_TITLE, 0, 0, 1, 1)
        self.GRID_PRES_CONTENT.addWidget(self.LNE_CNTT_TITLE, 0, 1, 1, 1)
        self.GRID_PRES_CONTENT.addItem(SPC_PRES_CNTT_V, 0, 2, 4, 1)
        self.GRID_PRES_CONTENT.addWidget(self.LBL_CNTT_SUBTITLE, 1, 0, 1, 1)
        self.GRID_PRES_CONTENT.addWidget(self.LNE_CNTT_SUBTITLE, 1, 1, 1, 1)
        self.GRID_PRES_CONTENT.addWidget(self.CHK_CNTT_DISPDATE, 2, 0, 1, 2)
        self.GRID_PRES_CONTENT.addWidget(self.CHK_CNTT_USEWIDESCR, 3, 0, 1, 2)

        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_TITLE, 0, 0, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.FCB_FAC_TITLE, 0, 1, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_TITLE_PREV, 0, 3, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.BTN_FAC_TITLE, 0, 4, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_SUBTITLE, 1, 0, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.FCB_FAC_SUBITITLE, 1, 1, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_SUBTITLE_PREV, 1, 3, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.BTN_FAC_SUBTITLE, 1, 4, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_BODY, 2, 0, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.FCB_FAC_BODY, 2, 1, 1, 1)
        self.GRID_PRES_FAC.addItem(SPC_PRES_FAC_H, 2, 2, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_BODY_PREV, 2, 3, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.BTN_FAC_BODY, 2, 4, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LIN_FAC_DIVIDER, 3, 0, 1, 5)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_ROLE, 4, 0, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.FCB_FAC_ROLE, 4, 1, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_ROLE_PREV, 4, 3, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.BTN_FAC_ROLE, 4, 4, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_NAME, 5, 0, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.FCB_FAC_NAME, 5, 1, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.LBL_FAC_NAME_PREV, 5, 3, 1, 1)
        self.GRID_PRES_FAC.addWidget(self.BTN_FAC_NAME, 5, 4, 1, 1)

        self.GRID_PRES_BG.addWidget(self.PIX_BG_PREVIEW, 0, 0, 2, 1)
        self.GRID_PRES_BG.addWidget(self.LBL_BG_PREVIEW, 2, 0, 1, 1)
        self.GRID_PRES_BG.addWidget(self.BTN_BG_BROWSE, 0, 1, 1, 1)
        self.GRID_PRES_BG.addWidget(self.BTN_BG_DISCARD, 1, 1, 1, 1)
        self.GRID_PRES_BG.addItem(SPC_PRES_BG_H, 2, 1, 1, 1)

        ## Members Tab
        self.GRID_MEMBERS.addWidget(self.LNE_MEM_SEARCHADD, 2, 0, 1, 1)
        self.GRID_MEMBERS.addWidget(self.LST_MEM_MEMBERS, 3, 0, 6, 1)
        self.GRID_MEMBERS.addWidget(self.BTN_MEM_ADD, 2, 1, 1, 1)
        self.GRID_MEMBERS.addWidget(self.BTN_MEM_EDIT, 3, 1, 1, 1)
        self.GRID_MEMBERS.addWidget(self.BTN_MEM_REMOVE, 4, 1, 1, 1)
        self.GRID_MEMBERS.addItem(SPC_MEM_V1, 5, 1, 1, 1)
        self.GRID_MEMBERS.addWidget(self.BTN_MEM_IMPORT, 6, 1, 1, 1)
        self.GRID_MEMBERS.addWidget(self.BTN_MEM_SAVE, 7, 1, 1, 1)
        self.GRID_MEMBERS.addWidget(self.BTN_MEM_EXPORT, 8, 1, 1, 1)
        self.GRID_MEMBERS.addWidget(self.GBX_MEM_DETAILS, 10, 0, 1, 2)

        self.GRID_MEM_DET.addWidget(self.LBL_DET_MEMBERS, 0, 0, 1, 1)
        self.GRID_MEM_DET.addWidget(self.LBL_DET_MEN, 0, 1, 1, 1)
        self.GRID_MEM_DET.addWidget(self.LBL_DET_OTHERS, 1, 0, 1, 1)
        self.GRID_MEM_DET.addWidget(self.LBL_DET_WOMEN, 1, 1, 1, 1)

        ## Initialization
        self.setupDisplay()
        self.setupValidators()
        MEM.setup()
        self.TBW_SETTINGS.setCurrentIndex(0)
        self.setupConnections()

    
    def setupValidators(self):
        self.LNE_MEM_SEARCHADD.setValidator(QtGui.QRegExpValidator(QtCore.QRegExp("[a-z-A-Z. -]+")))


    def setupDisplay(self):
        self.setWindowTitle("Settings")
        self.BTN_OK.setText("OK")
        self.LBL_APPVERSION.setText(f"{SW.NAME} v{SW.VERSION}")

        ## General Tab
        self.TBW_SETTINGS.setTabText(self.TBW_SETTINGS.indexOf(self.TAB_GENERAL), "General")
        self.GBX_GEN_APPEARANCE.setTitle("Appearance")
        self.GBX_GEN_PRESETS.setTitle("Presets")
        self.CHK_GEN_ALWAYS_ON_TOP.setText("Always on top")
        self.CHK_GEN_USE_SUGGESTED.setText("Automatically use suggested participant")
        self.CHK_CNTT_ENABLECBXSCROLL.setText("Change participant when scrolling through field")
        self.BTN_GEN_IMPORT.setToolTip("Import a preset")
        self.BTN_GEN_MODIFY.setToolTip("Modify this preset")
        self.BTN_GEN_REMOVE.setToolTip("Remove this preset from this list")

        PRESETS = [
                "Sabbath Service",
                "Adventist Youth Service",
                "Midweek Service",
                "Midweek Service",
                "Vesper Service",
                "District Fellowship"
                ]
        self.LST_GEN_PRESETS.addItems(PRESETS)
        
        ## Presentation Tab
        self.TBW_SETTINGS.setTabText(self.TBW_SETTINGS.indexOf(self.TAB_PRESENTATION), "Presentation")
        self.GBX_PRES_CONTENT.setTitle("Content")
        self.GBX_PRES_FONTSCOLORS.setTitle("Fonts and Colors")
        self.GBX_PRES_BACKGROUND.setTitle("Background")

        self.CHK_CNTT_DISPDATE.setText("Display date")
        self.CHK_CNTT_USEWIDESCR.setText("Use widescreen (16:9)")

        self.LBL_CNTT_TITLE.setText("Title:")
        self.LBL_CNTT_SUBTITLE.setText("Subtitle:")
        self.LNE_CNTT_TITLE.setPlaceholderText("Sabbath Service Participants")
        self.LNE_CNTT_SUBTITLE.setPlaceholderText("Happy Sabbath!")

        self.LBL_FAC_TITLE.setText("Title:")
        self.LBL_FAC_SUBTITLE.setText("Subtitle:")
        self.LBL_FAC_BODY.setText("Body:")
        self.LBL_FAC_NAME.setText("Name:")
        self.LBL_FAC_ROLE.setText("Role:")

        PREVIEWS = [self.LBL_FAC_TITLE_PREV, self.LBL_FAC_SUBTITLE_PREV, self.LBL_FAC_BODY_PREV, self.LBL_FAC_ROLE_PREV, self.LBL_FAC_NAME_PREV]
        for obj in PREVIEWS: obj.setText("Preview")
        self.LBL_BG_PREVIEW.setText("Preview")

        self.BTN_PRES_RESET.setText("Reset to defaults")
        
        ## Members Tab
        self.TBW_SETTINGS.setTabText(self.TBW_SETTINGS.indexOf(self.TAB_MEMBERS), "Members")
        self.LNE_MEM_SEARCHADD.setPlaceholderText("Search or type to add new member")
        self.BTN_MEM_ADD.setToolTip("Add this member to the list")
        self.BTN_MEM_EDIT.setToolTip("Rename this member")
        self.BTN_MEM_REMOVE.setToolTip("Remove this member from the list")
        self.BTN_MEM_SAVE.setToolTip("Save current member list")
        self.BTN_MEM_IMPORT.setToolTip("Import member list")
        self.BTN_MEM_EXPORT.setToolTip("Export member list")

        self.GBX_MEM_DETAILS.setTitle("Details")


    def setupConnections(self):
        """
        Manages Signals and Slots from user interactions
        """
        self.BTN_OK.clicked.connect(lambda: self.saveChanges())
        self.BTN_BG_BROWSE.clicked.connect(lambda: UIB.browseForBackgroundImage())
        self.BTN_BG_DISCARD.clicked.connect(lambda: UIB.discardImage())
        self.LNE_CNTT_TITLE.textChanged.connect(lambda: UIB.updatePackage())
        self.LNE_CNTT_SUBTITLE.textChanged.connect(lambda: UIB.updatePackage())

        ## Members
        self.LNE_MEM_SEARCHADD.mouseReleaseEvent = lambda event: MEM.searchClicked(event)
        self.LNE_MEM_SEARCHADD.textChanged.connect(lambda: MEM.checkSearchAdd())
        self.LST_MEM_MEMBERS.itemSelectionChanged.connect(lambda: MEM.itemChanged())
        self.BTN_MEM_ADD.clicked.connect(lambda: MEM.addNewMember())
        self.BTN_MEM_EDIT.clicked.connect(lambda: MEM.editMember())
        self.BTN_MEM_REMOVE.clicked.connect(lambda: MEM.removeMember())
        self.BTN_MEM_SAVE.clicked.connect(lambda: MEM.saveMemberList())
        self.BTN_MEM_IMPORT.clicked.connect(lambda: MEM.importMemberList())
        self.BTN_MEM_EXPORT.clicked.connect(lambda: MEM.exportMemberList())


        ## Exceptions
        self.CHK_CNTT_USEWIDESCR.setEnabled(False)
        self.CHK_GEN_ALWAYS_ON_TOP.setEnabled(False)


    def keyPressEvent(self, event):
        if event.key() == QtCore.Qt.Key_Escape:
            self.closeEvent()

        if event.key() in [QtCore.Qt.Key_Enter, QtCore.Qt.Key_Return]:
            ## Shortcut Key for Adding
            if self.LNE_MEM_SEARCHADD.hasFocus():
                MEM.addNewMember()


    def closeEvent(self, event=None):
        """
        Close event for Settings window.
        Triggers when the user selected "OK" or "x" button.
        """
        if SYS.GLOBAL_STATE == 3: return
        if self.CHANGED:
            MSG_BOX = QtWidgets.QMessageBox()
            MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON))
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Question)
            MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.No)
            MSG_BOX.setText("You have unsaved changes\nDo you want to save your preferences?")
            MSG_BOX.setWindowTitle(SW.NAME)
            MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
            MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')
            CORE.centerInsideWindow(MSG_BOX, UIB)
            
            RET = MSG_BOX.exec_()
            if RET == QtWidgets.QMessageBox.Yes:
                self.saveChanges()
            else:
                PKG.restart()
                self.hide()
        else:
            self.saveChanges()
    

    def saveChanges(self):
        """
        Triggers when user confirms and exits the settings window
        """
        DCFG["CONFIG"].update({"TXT_TITLE": PKG.TXT_TITLE})
        DCFG["CONFIG"].update({"TXT_SUBTITLE": PKG.TXT_SUBTITLE})
        PDB.dump()
        self.hide()


    def enterWindow(self):
        """
        Triggers after UIA's settings button was clicked

        Also handles dynamic position of Settings depending on UIA's pos
        """
        ## Default values
        self.ENTERING = True
        self.BTN_MEM_ADD.setEnabled(False)

        ## Spawn settings at center of UIA
        if self.isHidden():
            CORE.centerInsideWindow(self, UIA)
        
        ## Restore window / Show
        self.show()
        self.setWindowState(self.windowState() & ~QtCore.Qt.WindowMinimized | QtCore.Qt.WindowActive)
        self.activateWindow()

        ## Re-initialize variable
        self.LNE_CNTT_TITLE.setText(PKG.TXT_TITLE)
        self.LNE_CNTT_SUBTITLE.setText(PKG.TXT_SUBTITLE)
        self.ENTERING = False
        self.CHANGED = False
    

    def updatePackage(self):
        """
        Updates the current package to what's changed
        """
        if self.ENTERING: return
        self.CHANGED = True                                 ## Indicating that the settings window is modified
        PKG.TXT_TITLE = self.LNE_CNTT_TITLE.text()
        PKG.TXT_SUBTITLE = self.LNE_CNTT_SUBTITLE.text()


    def browseForBackgroundImage(self):
        """
        Customizes the background image
        """
        PATH_BGIMG = QtWidgets.QFileDialog.getOpenFileName(None, 'Browse for Hymnal Package', os.getcwd(),
                    f'Images ({" ".join(["*.{}".format(fo.data().decode()) for fo in QtGui.QImageReader.supportedImageFormats()])})'.format())
                    
        if PATH_BGIMG[0] != '':   
            self.PIX_BG_PREVIEW.setPixmap(QtGui.QPixmap(PATH_BGIMG[0]).scaledToHeight(85, QtCore.Qt.KeepAspectRatio & Qt.SmoothTransformation))
            PKG.IMG_BACKGROUND = PATH_BGIMG[0]
            self.updateBackgroundImage(PATH_BGIMG[0])


    def updateBackgroundImage(self, bg):
        """
        Makes copy of the imported bg image to program directory and updating the configuration
        """
        for i in ['png', 'jpg', 'svg', 'bmp', 'tiff', 'jpeg', 'gif']:
            try: os.remove(f"{SYS.DIR_PROGRAM}/BG.{i}")
            except FileNotFoundError: pass

        shutil.copy(bg, SYS.DIR_PROGRAM)
        FILE = bg.split('/')[-1]
        DEST = f"{SYS.DIR_PROGRAM}/BG{os.path.splitext(FILE)[1]}"

        try: os.rename(f"{SYS.DIR_PROGRAM}/{FILE}", DEST)
        except (FileNotFoundError, FileExistsError) as e: LOG.debug(e)
        DCFG["CONFIG"].update({"IMG_BACKGROUND": DEST}); PDB.dump()
    

    def discardImage(self):
        """
        Discards and restores the background image to default
        """
        PKG.IMG_BACKGROUND = PKG.DEF_IMG_BACKGROUND
        self.PIX_BG_PREVIEW.setPixmap(QtGui.QPixmap(PKG.IMG_BACKGROUND).scaledToHeight(self.BG_SCALING, QtCore.Qt.KeepAspectRatio & Qt.SmoothTransformation))
        self.updateBackgroundImage(PKG.IMG_BACKGROUND)




class Members(object):
    """
    Handles all member functions from Settings
    """
    def __init__(self):
        self.DUPLICATE_THRESHOLD = 0.85                                     ## Sets how sensitive is the duplicate scanner
        self.ITEMS_THRESHOLD = 7                                            ## Number of names to be allowed for exporting a member list
        self.SEARCH_STATE = False
        self.LAST_INPUT = ''
        self.SAVE_STATE = True


    def setup(self):
        UIB.LST_MEM_MEMBERS.clear()
        UIB.LST_MEM_MEMBERS.addItems(DCFG["POOL"]["NAMES"])
        self.CACHED_MEMBERS = self.getCurrentMembers()
        self.generalRefresh()


    def getCurrentMembers(self, lowercase=False):
        """
        Returns a string names of members from the list widget
        """
        return ([UIB.LST_MEM_MEMBERS.item(i).text() if not lowercase else
            UIB.LST_MEM_MEMBERS.item(i).text().lower() for i in range(UIB.LST_MEM_MEMBERS.count())])


    def checkSearchAdd(self):
        """
        Handles search/add bar (line edit)
        """
        UIB.BTN_MEM_ADD.setToolTip(f'Add {UIB.LNE_MEM_SEARCHADD.text()} to member list')


        if len(UIB.LNE_MEM_SEARCHADD.text()) < len(self.LAST_INPUT):
            self.SEARCH_STATE = True
            UIB.LST_MEM_MEMBERS.clear()
            UIB.LST_MEM_MEMBERS.addItems(self.CACHED_MEMBERS)

        if UIB.LNE_MEM_SEARCHADD.text() != '':
            UIB.BTN_MEM_ADD.setEnabled(True)
            self.SEARCH_STATE = True
            self.filterItems()

        if not len(re.sub(r"[^A-Za-z]+", '', UIB.LNE_MEM_SEARCHADD.text())):
            UIB.BTN_MEM_ADD.setEnabled(False)
            self.SEARCH_STATE = False
        
        self.generalRefresh()

    
    def filterItems(self):
        INPUT = re.sub(r"[^A-Za-z]+", '', UIB.LNE_MEM_SEARCHADD.text().lower())                             ## Use RegEx to filter out non-alphabet characters
        OUTPUT = [n for n in self.CACHED_MEMBERS if INPUT in re.sub(r"[^A-Za-z]+", '', n.lower())]
        UIB.LST_MEM_MEMBERS.clear()
        UIB.LST_MEM_MEMBERS.addItems(OUTPUT)
        UIB.LST_MEM_MEMBERS.sortItems()
        self.LAST_INPUT = UIB.LNE_MEM_SEARCHADD.text()
        self.refreshButtons()
    

    def hasDuplicateName(self, name:str, renaming=None):
        return True if len(self.getSimilarNames(name, renaming)) else False


    def getSimilarNames(self, name:str, renaming=None):    
        """
        Uses Levenshtein Ratio method to determine the similarity
        of the existing names vs the proposed name entry
        """
        MEMBERS = self.CACHED_MEMBERS
        if renaming is not None: MEMBERS.remove(renaming) 
        return [n for n in MEMBERS if lev.ratio(n.lower(), name.lower()) > self.DUPLICATE_THRESHOLD]
        

    def displayDialog(self, mode, similar:list=None, name=''):
        """
        Constructor method for similar names dialog
        """
        if mode == 0:
            LOG.warn(f"Members: Duplicate name detected: {name}")
            MSG_BOX = QtWidgets.QMessageBox(); MSG_BOX.setWindowTitle(SW.NAME)
            MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON)); MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Ok); MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.Ok)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Warning)
            MSG_BOX.setText(f"{name} is already in the list.")
            MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')

        elif mode == 1:
            LOG.info(f'Members: Similar names detected: {", ".join(similar)}')
            MSG_BOX = QtWidgets.QMessageBox(); MSG_BOX.setWindowTitle(SW.NAME)
            MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON)); MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Question)
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No); MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.No)
            MSG_BOX.setText("The name you entered might be already in the list.\nContinue anyway?")
            MSG_BOX.setDetailedText("Similar name(s):\n\n" + "\n".join(map(str, [f'{str(i+1).zfill(len(str(len(similar))))}. {n}' for i, n in enumerate(similar)])))
            MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')
           
        elif mode == 2:
            LOG.info(f'Members: Too few members for exporting')
            MSG_BOX = QtWidgets.QMessageBox(); MSG_BOX.setWindowTitle(SW.NAME)
            MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON)); MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
            MSG_BOX.setIcon(QtWidgets.QMessageBox.Question)
            MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Yes | QtWidgets.QMessageBox.No); MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.No)
            MSG_BOX.setText(f"The list only contains {len(self.getCurrentMembers())} member(s)\nDo you still want to continue?")
            MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')

        return MSG_BOX


    def adjustNameFormat(self, name:str):
        """
        Reformats the name string into a presentable name:

        • Auto-capitalizes name
        • Auto-fill period after "Bro" or "Sis"
        """
        NAME = ' '.join(f"{c[0].upper()}{c[1:].lower()}" for c in name.split())                 ## Convert into Name Case which auto-capitalizes first letter of every word
        try:
            if NAME.split()[0].lower()[:3] in ['bro', 'sis', 'ptr'] and NAME.split()[0].lower() not in ['brother', 'sister', 'pastor']:
                if NAME.split()[0][-1] != '.':
                    NAME = f"{NAME.split()[0]}. {' '.join(map(str, NAME.split()[1:]))}"
        except IndexError:                                                                      ## Handler for blank field
            pass
        return NAME


    def addNewMember(self):
        """
        Inserts new item using the text from search bar
        Also checks for duplicate and similar member names
        """
        if not UIB.BTN_MEM_ADD.isEnabled(): return
        ## Check for duplicate
        PROPOSED = UIB.LNE_MEM_SEARCHADD.text().strip()
        if PROPOSED.lower() in map(lambda i: i.lower().strip(), self.CACHED_MEMBERS):
            self.displayDialog(0, self.getSimilarNames(PROPOSED), PROPOSED).exec_(); return

        elif self.hasDuplicateName(PROPOSED):
            if self.displayDialog(1, self.getSimilarNames(PROPOSED)).exec_() != QtWidgets.QMessageBox.Yes: return

        UIB.LNE_MEM_SEARCHADD.clear()
        UIB.LST_MEM_MEMBERS.addItem(self.adjustNameFormat(PROPOSED))
        UIB.LST_MEM_MEMBERS.sortItems()
        self.CACHED_MEMBERS = self.getCurrentMembers()
        UIB.LST_MEM_MEMBERS.findItems(self.adjustNameFormat(PROPOSED), QtCore.Qt.MatchExactly)[0].setSelected(True)
        self.SAVE_STATE = False
        self.generalRefresh()


    def editMember(self):
        """
        Lets the user edit and rename member
        """
        OLD_NAME = UIB.LST_MEM_MEMBERS.selectedItems()[0].text()

        class DLG_RENAME(QtWidgets.QDialog):
            def __init__(self, parent = None):
                QtWidgets.QWidget.__init__(self, parent)
                self.setWindowTitle(SW.NAME)
                self.setFixedSize(270, 110)
                self.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON))
                self.setWindowFlags(self.windowFlags() | QtCore.Qt.WindowStaysOnTopHint)
                self.GRID_MAIN = QtWidgets.QGridLayout(self); self.GRID_MAIN.setObjectName("GRID_MAIN")
                self.LBL_RENAME = QtWidgets.QLabel(self); self.LBL_RENAME.setObjectName("LBL_RENAME")
                SPC_V1 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
                SPC_V2 = QtWidgets.QSpacerItem(20, 40, QtWidgets.QSizePolicy.Minimum, QtWidgets.QSizePolicy.Expanding)
                self.LNE_RENAME = QtWidgets.QLineEdit(self); self.LNE_RENAME.setObjectName("LNE_RENAME")
                self.LNE_RENAME.setValidator(QtGui.QRegExpValidator(QtCore.QRegExp("[a-z-A-Z. -]+")))
                self.BTN_BOX = QtWidgets.QDialogButtonBox(self); self.BTN_BOX.setObjectName("BTN_BOX")
                self.BTN_BOX.setStandardButtons(QtWidgets.QDialogButtonBox.Cancel|QtWidgets.QDialogButtonBox.Ok)

                self.GRID_MAIN.addWidget(self.LBL_RENAME, 1, 0, 1, 1)
                self.GRID_MAIN.addItem(SPC_V1, 3, 0, 1, 1)
                self.GRID_MAIN.addItem(SPC_V2, 0, 0, 1, 1)
                self.GRID_MAIN.addWidget(self.LNE_RENAME, 2, 0, 1, 1)
                self.GRID_MAIN.addWidget(self.BTN_BOX, 5, 0, 1, 1)

                self.LBL_RENAME.setText(f"Set a new name for {OLD_NAME}")
                self.LNE_RENAME.setText(OLD_NAME)
                self.BTN_BOX.accepted.connect(lambda: self.accept())
                self.BTN_BOX.rejected.connect(lambda: self.reject())
                self.setStyleSheet('QPushButton {min-width: 50px;}')
                CORE.centerInsideWindow(self, UIB)
                self.exec_()
            
            def accept(self):
                NEW_NAME = self.LNE_RENAME.text().strip()
                if NEW_NAME != '':
                    if NEW_NAME.lower() == OLD_NAME.lower():
                        self.hide()

                    elif NEW_NAME.lower() in MEM.getCurrentMembers(True):
                        MEM.displayDialog(0, MEM.getSimilarNames(NEW_NAME), NEW_NAME).exec_(); return

                    elif MEM.hasDuplicateName(NEW_NAME, OLD_NAME):
                        if MEM.displayDialog(1, MEM.getSimilarNames(NEW_NAME)).exec_() != QtWidgets.QMessageBox.Yes: return     

                    UIB.LST_MEM_MEMBERS.selectedItems()[0].setText(MEM.adjustNameFormat(NEW_NAME))
                    MEM.CACHED_MEMBERS = MEM.getCurrentMembers()         
                self.hide()      

            def reject(self):
                self.hide()
        
        gc.collect()
        DLG_RENAME()
        self.SAVE_STATE = False
        self.generalRefresh()
            
    
    def removeMember(self):
        for i in UIB.LST_MEM_MEMBERS.selectedItems():
            UIB.LST_MEM_MEMBERS.takeItem(UIB.LST_MEM_MEMBERS.row(i))

        self.CACHED_MEMBERS = self.getCurrentMembers()
        self.SAVE_STATE = False
        self.generalRefresh()


    def itemChanged(self):
        UIB.BTN_MEM_ADD.setEnabled(False)

        self.generalRefresh() 

        try:
            if len(UIB.LST_MEM_MEMBERS.selectedItems()):
                UIB.BTN_MEM_EDIT.setToolTip(f'Rename "{UIB.LST_MEM_MEMBERS.currentItem().text()}"')
                UIB.BTN_MEM_REMOVE.setToolTip(f'Remove "{UIB.LST_MEM_MEMBERS.currentItem().text()}" from the list')
        except AttributeError:
            UIB.BTN_MEM_EDIT.setToolTip(f'Rename this member')
            UIB.BTN_MEM_REMOVE.setToolTip(f'Remove this member from the list')


    def searchClicked(self, event):
        UIB.BTN_MEM_ADD.setEnabled(True if len(self.CACHED_MEMBERS) and len(re.sub(r"[^A-Za-z]+", '', UIB.LNE_MEM_SEARCHADD.text())) else False) 


    def generalRefresh(self):
        """
        Works as a general refresher for multiple functions
        """
        self.refreshButtons()
        self.refreshDetails()
        UIB.LST_MEM_MEMBERS.sortItems()


    def refreshButtons(self):
        """
        Handles all button-related items for Members tab
        """
        if len(UIB.LST_MEM_MEMBERS.selectedItems()):
            UIB.BTN_MEM_REMOVE.setEnabled(True)
            UIB.BTN_MEM_EDIT.setEnabled(True)
        else:
            UIB.BTN_MEM_REMOVE.setEnabled(False); UIB.BTN_MEM_REMOVE.setToolTip('')
            UIB.BTN_MEM_EDIT.setEnabled(False); UIB.BTN_MEM_EDIT.setToolTip('')

        UIB.BTN_MEM_EXPORT.setEnabled(True if len(self.getCurrentMembers()) and not self.SEARCH_STATE else False)
        UIB.BTN_MEM_SAVE.setEnabled(True if len(self.getCurrentMembers()) and not self.SAVE_STATE else False)
    

    def refreshDetails(self):
        """
        Refreshes the details tab 
        """
        ITEMS = self.CACHED_MEMBERS
        DATA = [len(self.CACHED_MEMBERS), 0, 0, 0]                                                                   ## Members, Men, Women, Others

        DATA[1] = len(list(filter(re.compile("(?i)^bro[. ]|brother").match, ITEMS)))
        DATA[2] = len(list(filter(re.compile("(?i)^sis[. ]|sister").match, ITEMS)))
        DATA[3] = DATA[0] - (DATA[1] + DATA[2])

        UIB.LBL_DET_MEMBERS.setEnabled(True if DATA[0] else False)
        UIB.LBL_DET_MEN.setEnabled(True if DATA[1] else False)
        UIB.LBL_DET_WOMEN.setEnabled(True if DATA[2] else False)
        UIB.LBL_DET_OTHERS.setEnabled(True if DATA[3] else False)

        if not DATA[0]:
            UIB.LBL_DET_MEMBERS.setText('There are no members here...')
            UIB.LNE_MEM_SEARCHADD.setPlaceholderText('Type the name of your first member here')
            UIB.LBL_DET_MEN.clear()
            UIB.LBL_DET_WOMEN.clear()
            UIB.LBL_DET_OTHERS.clear()
            return

        UIB.LNE_MEM_SEARCHADD.setPlaceholderText('Search or type to add new member')
        UIB.LBL_DET_MEMBERS.setText(f"Members: <b>{DATA[0]}</b>")

        PLOT = {
            UIB.LBL_DET_MEN: ("Men", DATA[1]),
            UIB.LBL_DET_WOMEN: ("Women", DATA[2]),
            UIB.LBL_DET_OTHERS: ("Others", DATA[3])
        }
        for k,v in PLOT.items():
            k.setText(f'{v[0]}: ' + (f"<b>{v[1]}</b> <font color={QSS.TXT_DISABLED}>({round((v[1]*100)/DATA[0])}%)</font>" if v[1] else '...'))
        

    def importMemberList(self):
        FILE, EXT = QtWidgets.QFileDialog.getOpenFileName(None, "Import Member List", PKG.DIR_IMPORT_MEMLIST, f'*.{SYS.EXT_MEMLIST}')
        if FILE:
            try:
                with open(FILE, 'rb') as f:
                    LOAD = json.loads(bz2.decompress(f.read()).decode('utf-8'))
            except (OSError, json.decoder.JSONDecodeError) as e:
                LOG.error(f'Invalid member list file: {FILE}')
                MSG_BOX = QtWidgets.QMessageBox(); MSG_BOX.setWindowTitle(SW.NAME)
                MSG_BOX.setWindowIcon(QtGui.QIcon(SYS.RES_APP_ICON)); MSG_BOX.setWindowFlags(Qt.WindowStaysOnTopHint)
                MSG_BOX.setStandardButtons(QtWidgets.QMessageBox.Ok); MSG_BOX.setDefaultButton(QtWidgets.QMessageBox.Ok)
                MSG_BOX.setIcon(QtWidgets.QMessageBox.Warning)
                MSG_BOX.setText(f"File is corrupted or invalid.")
                MSG_BOX.setStyleSheet('QPushButton {min-width: 50px;}')
                MSG_BOX.exec_()
            else:
                ## Save to JSON
                DCFG['CONFIG'].update({"DIR_IMPORT_MEMLIST": FILE})
                DCFG['POOL'].update({"NAMES": LOAD['MEMBER_LIST']})
                self.setup()
                PKG.DIR_IMPORT_MEMLIST = FILE
                PDB.dump()
                LOG.info(f"Member list \"{FILE}\" was successfully loaded.")


    def exportMemberList(self):
        """
        Export current values into a encrypted JSON file
        """
        if len(self.CACHED_MEMBERS) < self.ITEMS_THRESHOLD:
            if self.displayDialog(2).exec_() != QtWidgets.QMessageBox.Yes: return

        MEMBERS = {"MEMBER_LIST": self.CACHED_MEMBERS, "DATE_GENERATED": time.time()}
        MEMBERS = json.dumps(MEMBERS, indent=None, sort_keys=True)
        
        FILE, EXT = QtWidgets.QFileDialog.getSaveFileName(None, "Export Member List", PKG.DIR_EXPORT_MEMLIST, "*.prt")
        if FILE:
            LOG.info(f"Saving exported member list to {FILE}")
            try:
                with open(FILE, 'wb') as w:
                    w.write(bz2.compress(MEMBERS.encode('utf-8'), 9))
            except Exception as e:
                LOG.error(e)
            
            ## Save to JSON
            DCFG['CONFIG'].update({"DIR_EXPORT_MEMLIST": FILE})
            PKG.DIR_EXPORT_MEMLIST = FILE
            PDB.dump()
            LOG.info(f"Member list \"{FILE}\" was successfully exported.")
            self.saveMemberList(True)

    
    def saveMemberList(self, bypassDialog=False):
        """
        Saves the current member list
        """
        if len(self.CACHED_MEMBERS) < self.ITEMS_THRESHOLD and not bypassDialog:
            if self.displayDialog(2).exec_() != QtWidgets.QMessageBox.Yes: return
            
        DCFG['POOL'].update({"NAMES": self.CACHED_MEMBERS})
        PDB.dump()
        self.SAVE_STATE = True
        LOG.info(f"Member list was successfully saved.")
        self.generalRefresh()



class Settings(object):
    """
    Functional class for Settings window
    """
    def __init__(self):
        pass
    
    ## Some code from QWGT_SETTINGS will migrate here




if __name__ == '__main__':
    INIT_TIME = time.time()
    
    ## Create Application
    APP = QtWidgets.QApplication(sys.argv)

    ## Primary Software Initialization & Logging
    SW = KSoftware("Participants", "1.0.8", "Ken Verdadero, Reynald Ycong", file=__file__, parentName="MSDAC Systems", prodYear=2022, versionName="Release")
    LOG = KLog(System().DIR_LOG, __file__, SW.LOG_NAME_DATE(), SW.PY_NAME, SW.AUTHOR, cont=True, tms=True, delete_existing=True, tmsformat="%H:%M:%S.%f %m/%d/%y")

    SYS = System()
    SYS.verifyDirectories()

    PDB = Data()
    global DCFG
    DCFG = PDB.DATA
    RLS = DCFG['POOL']['ROLES']
    NMS = DCFG['POOL']['NAMES']

    LOG.info('Initializing Internal Classes')
    QSS = Stylesheet()
    PKG = Package()
    FLD = Fields()
    CORE = Core()
    EXP = Export()
    FMN = FileManager()
    STT = Settings()
    MEM = Members()

    LOG.info('Initializing UI')
    UIA = QWGT_PARTICIPANTS()
    UIB = QWGT_SETTINGS()
    
    UIA.setupUI()
    UIB.setupUI()

    UIA.closeEvent = lambda event: SYS.closeEvent(event)

    SYS.verifyRequisites()                                                                          ## Check for MS Office PowerPoint availability
    SYS.checkInstances() 

    ## Shows User Interface
    gc.collect()
    UIA.show()
    CORE.centerWindow(UIA)
    UIA.raise_()
    UIA.activateWindow()

    SYS.STARTUP_TIME = time.time()-INIT_TIME
    LOG.info(f'Initialization completed in {round(SYS.STARTUP_TIME, 3)} seconds.')
    sys.exit(APP.exec_())
    


